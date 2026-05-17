"""
build_p1.py — OpenAgentPay × AgentCore Payments 周会分享 PPT (Slide 1-7)

按 deep-research-ppt skill 规范：
  - 配色：金融科技主题（深蓝 #0A1A2E + 青 #00D4FF + 金 #FFB800）
    + AWS 橙（#FF9900）作 AgentCore 章节副色
    + 绿（#00FF88）作 OpenAgentPay 章节副色
  - 16:9 宽屏（13.33 × 7.5 英寸 = 标准 widescreen）
  - 章节分隔页 + 大数字指标 + 圆角卡片
  - 每页 200-500 字演讲备注（含 demo 演示指引）

Slides 1-7:
  1. 封面
  2. 目录
  3. 章节分隔页（第一章：什么是 AgentCore Payments）
  4. 场景钩子（你的 AI Agent 能自己付款吗？）
  5. AgentCore Payments 一句话定义 + 平台地位（图 01）
  6. x402 协议怎么工作（图 02）
  7. 经济学优势 — 3000x 差距（图 08）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================================
#  CONFIG
# ============================================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色 — 金融科技主题（深蓝 + 青 + 金）+ 章节副色（AWS 橙 / OAP 绿 / 警示红）
BG       = RGBColor(0x0A, 0x1A, 0x2E)   # 深蓝底
BG_CARD  = RGBColor(0x14, 0x2A, 0x44)   # 卡片色
BG_DARK  = RGBColor(0x05, 0x10, 0x1F)   # 章节分隔深底
FG       = RGBColor(0xF0, 0xF4, 0xF8)   # 主文字
DIM      = RGBColor(0x8B, 0xA0, 0xBA)   # 副文字
MUTED    = RGBColor(0x4A, 0x5C, 0x75)   # 辅助
CYAN     = RGBColor(0x00, 0xD4, 0xFF)   # 主色：青
GOLD     = RGBColor(0xFF, 0xB8, 0x00)   # 强调：金
AWS_ORG  = RGBColor(0xFF, 0x99, 0x00)   # AWS 章节
OAP_GRN  = RGBColor(0x00, 0xFF, 0x88)   # OpenAgentPay 章节
RED      = RGBColor(0xFF, 0x44, 0x44)   # 警示
PURPLE   = RGBColor(0xBF, 0x5A, 0xF2)   # 紫（次要）

IMG_DIR = Path(__file__).parent / "images"

# ============================================================================
#  HELPER FUNCTIONS (内联，避免 import 问题)
# ============================================================================
def bg(slide, color=BG):
    """设置整页背景色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def t(slide, left, top, w, h, text, *, size=18, color=FG, bold=False, align="left", font="Calibri"):
    """添加文本框"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb

def mt(slide, left, top, w, h, lines, font="Calibri"):
    """多行文本：lines 是 [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb

def bar(slide, left, top, w, h, color=CYAN):
    """装饰线条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def card(slide, left, top, w, h, color=BG_CARD, line_color=None):
    """圆角卡片"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def hdr(slide, title, subtitle=None, accent=CYAN):
    """页面标题栏 + 分隔线（顶部）"""
    t(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
      title, size=28, color=FG, bold=True)
    if subtitle:
        t(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
          subtitle, size=14, color=DIM)
    bar(slide, Inches(0.6), Inches(1.45), Inches(2), Emu(40000), accent)

def ftr(slide, page_num=None, total=14):
    """页脚（主题 + 页码）"""
    t(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
      "OpenAgentPay × AgentCore Payments  ·  2026-05-18  ·  Neo Sun",
      size=9, color=MUTED)
    if page_num:
        t(slide, Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.3),
          f"{page_num:02d} / {total:02d}", size=9, color=MUTED, align="right")

def n(slide, text):
    """设置演讲备注"""
    notes = slide.notes_slide.notes_text_frame
    notes.text = text

def metric(slide, left, top, w, h, number, label, color=CYAN, num_size=44, lbl_size=11):
    """大数字指标卡"""
    card(slide, left, top, w, h, BG_CARD, color)
    # 大数字
    t(slide, left, top + Inches(0.3), w, Inches(1.2),
      number, size=num_size, color=color, bold=True, align="center")
    # 标签
    t(slide, left, top + h - Inches(0.5), w, Inches(0.3),
      label, size=lbl_size, color=DIM, align="center")

def chap(prs, num, title, subtitle, accent=CYAN):
    """章节分隔页（深底大字）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG_DARK)
    # "第 N 章"
    t(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
      f"第 {num} 章", size=24, color=accent, bold=True)
    # 大标题
    t(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.5),
      title, size=56, color=FG, bold=True)
    # 副标题
    t(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.6),
      subtitle, size=18, color=DIM)
    # 装饰线
    bar(s, Inches(0.6), Inches(5.5), Inches(3), Emu(60000), accent)
    return s

# ============================================================================
#  BUILD SLIDES 1-7
# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================================
# Slide 1 — 封面
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_DARK)

# 顶部标签
t(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
  "AWS Web3 Team · 周会分享 · 2026-05-18", size=12, color=DIM)

# 大标题
t(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
  "AgentCore Payments", size=56, color=CYAN, bold=True)
t(s, Inches(0.6), Inches(3.4), Inches(12), Inches(1.2),
  "+ OpenAgentPay", size=56, color=GOLD, bold=True)

# 副标题
t(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
  "给 AI Agent 一张钱包 — 从 AWS 原版到亚洲生态扩展", size=20, color=FG)
t(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
  "Neo Sun  ·  AWS Solutions Architect  ·  jiasunm@amazon.com", size=14, color=DIM)

# 装饰线
bar(s, Inches(0.6), Inches(6.0), Inches(2), Emu(60000), CYAN)
bar(s, Inches(2.7), Inches(6.0), Inches(2), Emu(60000), GOLD)

n(s, """【封面】各位同事，今天周会我想分享一个我研究了两周的话题——Amazon Bedrock AgentCore Payments，以及我基于它做的一个扩展项目，叫 OpenAgentPay。

先给一个上下文：上周我发表了一篇公众号文章解读 AgentCore Payments，老板看了之后让我在团队周会上分享。但我觉得只讲 AgentCore Payments 太单一了——因为这个产品当前 Preview 阶段只支持 Coinbase 和 Stripe 两个钱包，我们 Web3 团队服务的所有亚洲客户、所有合规交易所、所有中心化钱包用户——全部用不上。

所以我今天的分享分两部分：前半讲 AgentCore Payments 是什么、做对了什么、有什么短板；后半讲我做的 OpenAgentPay 怎么解决这个短板，让我们的客户也能享受 AgentCore Payments 的能力。

整个演讲大概 15-20 分钟，欢迎大家随时打断提问。这个项目目前是 work-in-progress 状态——我会诚实地讲哪里完成了、哪里还在开发，哪里需要大家一起加入。""")
ftr(s, 1)

# ============================================================================
# Slide 2 — 目录
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "今天分享的两件事", "前半部分：理解 AgentCore Payments  ·  后半部分：OpenAgentPay 扩展方案")

# 左半 - Part 1
card(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5), BG_CARD, AWS_ORG)
t(s, Inches(1.1), Inches(2.3), Inches(5), Inches(0.5),
  "Part 1", size=16, color=AWS_ORG, bold=True)
t(s, Inches(1.1), Inches(2.7), Inches(5), Inches(0.6),
  "AgentCore Payments", size=28, color=FG, bold=True)
mt(s, Inches(1.1), Inches(3.5), Inches(5.4), Inches(2.8), [
    ("第 1 章  ·  什么是 AgentCore Payments", 14, FG, True),
    ("        场景钩子 + 平台定位 + x402 协议 + 经济学", 12, DIM, False),
    ("", 6, DIM, False),
    ("第 2 章  ·  现状与短板", 14, FG, True),
    ("        当前 Preview 数据 + Coinbase/Stripe 限制", 12, DIM, False),
])

# 右半 - Part 2
card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5), BG_CARD, OAP_GRN)
t(s, Inches(7.1), Inches(2.3), Inches(5), Inches(0.5),
  "Part 2", size=16, color=OAP_GRN, bold=True)
t(s, Inches(7.1), Inches(2.7), Inches(5), Inches(0.6),
  "OpenAgentPay", size=28, color=FG, bold=True)
mt(s, Inches(7.1), Inches(3.5), Inches(5.4), Inches(2.8), [
    ("第 3 章  ·  解决方案设计", 14, FG, True),
    ("        5 层架构 + 双协议轨道 + Live AWS demo", 12, DIM, False),
    ("", 6, DIM, False),
    ("第 4 章  ·  HashKey 实战 + 路线图", 14, FG, True),
    ("        链上 e2e demo + 三种集成形态 + CTA", 12, DIM, False),
])

n(s, """【目录】今天的分享分两部分。

第一部分讲 AgentCore Payments 本身——这是 AWS 上周发布的产品，我会用 5 页 slide 把它讲清楚：什么是它、它怎么工作（x402 协议）、它的经济学优势、它的现状、它的短板。

我特别想强调'短板'——不是因为这个产品不好，恰恰相反，这是个非常 promising 的产品。但它的支持范围有局限性，对我们 Web3 团队覆盖的亚洲客户来说，**当前不可用**。这就是第二部分的起点。

第二部分讲 OpenAgentPay——我做的扩展层。我会讲：它的设计思路（5 层可插拔架构）、它如何解决短板（双协议轨道）、它真的跑得通吗（HashKey Chain 链上 demo + AWS 部署）、未来怎么走（三种集成形态 + 社区化路线）。

整个演讲大概 15-20 分钟。我建议大家：技术问题随时打断，业务问题留到最后讨论环节。""")
ftr(s, 2)

# ============================================================================
# Slide 3 — 章节分隔页：第一章
# ============================================================================
s = chap(prs, 1, "什么是 AgentCore Payments", "AWS 给 AI Agent 发了一张钱包", AWS_ORG)
n(s, """【第一章】我们先看 AgentCore Payments 本身。这是 AWS 在 2026 年 5 月 7 日金融服务峰会上低调发布的一个产品，但它会被未来定义为'Agent 经济元年的里程碑'——因为它是首个面向自主 AI Agent 的托管支付基础设施。

接下来 5 页，我从一个场景钩子开始，讲清楚它是什么、它怎么工作、它的经济学，最后引出它的现状和短板。""")
ftr(s, 3)

# ============================================================================
# Slide 4 — 场景钩子
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)

# 顶部小标签
t(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
  "想象一下这个场景", size=12, color=DIM)

# 大问题
t(s, Inches(0.6), Inches(1.6), Inches(12), Inches(2.0),
  "你的 AI Agent", size=42, color=FG, bold=True)
t(s, Inches(0.6), Inches(2.5), Inches(12), Inches(2.0),
  "能自己付款吗？", size=72, color=GOLD, bold=True)

# 副文
t(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
  "在今天之前，答案是：不能。", size=24, color=DIM)
t(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
  "Agent 必须停下来等人类介入：开账号、绑信用卡、管理 API Key、处理订阅", size=14, color=MUTED)

# 底部金句
card(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6), BG_CARD, RED)
t(s, Inches(0.8), Inches(6.05), Inches(12), Inches(0.5),
  "智能是自主的。付款，不是。  ←  AI Agent 领域被所有人忽视但致命的矛盾",
  size=14, color=RED, bold=True)

n(s, """【场景钩子】我先问大家一个问题——你的 AI Agent 能自己付款吗？

想象一下这个场景：你的 AI Agent 在帮你做金融研究，分析到一半发现需要访问彭博终端的付费数据才能给出更准确的结论。这个时候 Agent 能自己付款吗？

在今天之前，答案是不能。Agent 必须停下来等人类介入：开账号、绑信用卡、管理 API Key、处理订阅...走完一整套人类世界为人类设计的支付流程。

这就是 AI Agent 领域一个被所有人忽视但致命的矛盾——智能是自主的，付款不是。

为什么这是个大问题？因为 Agent 经济正在快速到来。Anthropic 的 Computer Use、OpenAI 的 Operator、AWS 自己的 Strands Agent，都在让 Agent 能自主使用工具。但是只要付款这一环还要人参与，Agent 的自主性就被打了折扣。

AWS 的 AgentCore Payments 就是要解决这个问题——给 AI Agent 一张能自己用的钱包。这是 Agent 经济的关键基础设施层。""")
ftr(s, 4)

# ============================================================================
# Slide 5 — AgentCore Payments 一句话定义 + 平台地位
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "AgentCore Payments 是什么",
    "Amazon Bedrock AgentCore 平台的原生支付模块  ·  2026-05-07 Preview 发布", AWS_ORG)

# 一句话定义
card(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.0), BG_CARD, AWS_ORG)
t(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.4),
  "首个面向自主 AI Agent 的托管支付基础设施", size=18, color=GOLD, bold=True)
t(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(0.4),
  "让 Agent 在推理循环内自主发现 → 认证 → 结算 → 获取付费资源，无需人类介入",
  size=12, color=DIM)

# 嵌入 platform-landscape 图
img_path = str(IMG_DIR / "01-platform-landscape.png")
s.shapes.add_picture(img_path, Inches(2.0), Inches(2.9), Inches(9.3), Inches(3.7))

# 底部三个关键词
t(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
  "三个关键词：首个 (微软/Google 在 Agent 平台层都没同类)  ·  托管 (AWS 一手包办)  ·  推理循环内 (不打断 Agent 思考)",
  size=12, color=CYAN, align="center")

n(s, """【AgentCore Payments 是什么】2026 年 5 月 7 日，AWS 在金融服务峰会发布 Amazon Bedrock AgentCore Payments Preview。一句话定义：这是首个面向自主 AI Agent 的托管支付基础设施，让 Agent 在推理循环内自主发现、认证、结算、获取付费资源，无需人类介入。

注意三个关键词：

第一，首个。微软 Azure 和 Google Cloud 在 Agent 平台层都**没有**同类产品。AWS 领先了整整一代——这是 AWS 在 Agent 经济的战略卡位。

第二，托管。AWS 一手包办协议协商、钱包管理、限额控制、合规检查、可观测性。开发者只需要几行代码就能接入，不用自己写复杂的钱包逻辑。

第三，推理循环内。这是最关键的设计——付款动作不中断 Agent 思考，就像人喘气一样自然。Agent 调用工具的时候碰到 402 Payment Required，AgentCore 自动处理付款，Agent 完全无感。

看这张图，AgentCore Payments 不是独立产品，而是 AgentCore 平台的原生模块，跟 Identity、Gateway、Observability、Memory、Runtime 平起平坐。这就是 AWS 的设计哲学——'native, not bolted-on'，原生集成不是贴片。意味着企业安全团队审批一次就行，不用重新过评审。

它包含 4 个子能力：Payment Manager（编排大脑）、Payment Guardrail（风控边界）、Wallet Providers（钱包对接）、Protocols（协议层）。下一页讲它的核心协议 x402 怎么工作。""")
ftr(s, 5)

# ============================================================================
# Slide 6 — x402 协议怎么工作
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "x402 协议  ·  把 HTTP 402 状态码激活了",
    "Coinbase 2024 年底开放协议  ·  5 步搞定一次 Agent 付款", AWS_ORG)

# 嵌入 x402-sequence 图
img_path = str(IMG_DIR / "02-x402-sequence.png")
s.shapes.add_picture(img_path, Inches(2.0), Inches(1.7), Inches(9.3), Inches(3.7))

# 底部数据卡片：5 个 metric
metric(s, Inches(0.6),  Inches(5.7), Inches(2.4), Inches(1.2),
       "75M", "近 30 天 tx", CYAN, num_size=32)
metric(s, Inches(3.1),  Inches(5.7), Inches(2.4), Inches(1.2),
       "169M+", "累计 tx", GOLD, num_size=32)
metric(s, Inches(5.6),  Inches(5.7), Inches(2.4), Inches(1.2),
       "590K+", "活跃买家", OAP_GRN, num_size=32)
metric(s, Inches(8.1),  Inches(5.7), Inches(2.4), Inches(1.2),
       "100K+", "活跃卖家", PURPLE, num_size=32)
metric(s, Inches(10.6), Inches(5.7), Inches(2.1), Inches(1.2),
       "~2s", "L2 finalize", AWS_ORG, num_size=32)

n(s, """【x402 协议】要理解 AgentCore Payments，必须先懂 x402 协议。x402 是 Coinbase 在 2024 年底发布的开放协议，做了一件精妙的事——把 HTTP 协议里那个一直没人用的状态码 402 Payment Required 激活了。

这张时序图告诉大家 5 步搞定一次 Agent 付款：

第 1 步：Agent 发起请求 GET /resource，普通 HTTP 调用。
第 2 步：服务端返回 402 Payment Required，告诉你要付多少、付给谁、用啥币。
第 3 步：Agent 内部，AgentCore Payments 查预算 → 拿钱包私钥 → 做 EIP-3009 离线签名。
第 4 步：Agent 带上 X-Payment 头重试请求。
第 5 步：服务端验证签名 → Facilitator 上链结算 → 返回 200 OK + 你要的内容。

整个流程，**Agent 不需要持有 gas、不需要连区块链节点、不管理 nonce**。复杂度全部被 x402 Facilitator 吸收掉了，Agent 的 UX 跟调 REST API 一模一样。

底部这一组数字告诉大家 x402 不是玩具：30 天 7541 万笔交易、累计 1.69 亿笔、活跃买家 59 万、卖家 10 万、L2 结算约 2 秒完成 finalize（Flashblock 预确认 200 毫秒）。

被 Stripe、AWS、Cloudflare、Vercel、Quicknode、Messari、Alchemy 都信赖。AWS 选 x402 不是保守的选择——是押**最成熟、已生产验证、增长最快**的赛道。这是个深思熟虑的战略选型。""")
ftr(s, 6)

# ============================================================================
# Slide 7 — 经济学优势
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "为什么不能用传统信用卡  ·  经济学不允许",
    "x402 vs Stripe 信用卡  ·  扎心的 3000 倍差距", AWS_ORG)

# 左半：图 08
img_path = str(IMG_DIR / "08-economics.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(7.0), Inches(4.3))

# 右半：扎心数据卡片
card(s, Inches(7.9), Inches(1.7), Inches(4.8), Inches(2.0), BG_CARD, RED)
t(s, Inches(8.1), Inches(1.85), Inches(4.5), Inches(0.4),
  "传统 Stripe 信用卡", size=14, color=RED, bold=True)
t(s, Inches(8.1), Inches(2.3), Inches(4.5), Inches(0.8),
  "$0.30", size=44, color=RED, bold=True)
t(s, Inches(8.1), Inches(3.1), Inches(4.5), Inches(0.5),
  "单笔最低手续费", size=12, color=DIM)

card(s, Inches(7.9), Inches(3.85), Inches(4.8), Inches(2.0), BG_CARD, OAP_GRN)
t(s, Inches(8.1), Inches(4.0), Inches(4.5), Inches(0.4),
  "x402 + USDC on Base L2", size=14, color=OAP_GRN, bold=True)
t(s, Inches(8.1), Inches(4.45), Inches(4.5), Inches(0.8),
  "$0.0001", size=44, color=OAP_GRN, bold=True)
t(s, Inches(8.1), Inches(5.25), Inches(4.5), Inches(0.5),
  "单笔成本", size=12, color=DIM)

# 底部金句
card(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7), BG_CARD, GOLD)
t(s, Inches(0.8), Inches(6.18), Inches(12), Inches(0.55),
  "差了 3000 倍。这就是为什么所有按次计费的 Pay-per-use 市场都上不了规模 — 直到 x402 出现。",
  size=14, color=GOLD, bold=True, align="center")

n(s, """【经济学】你可能问：让 Agent 自动付款，直接绑张信用卡不就完了？答案是：经济学不允许。

看这张图：传统 Stripe 信用卡单笔最低手续费 0.30 美元，x402 单笔成本约 0.0001 美元——差了整整 3000 倍。

这意味着什么？大家想想：Agent 调用一次 API 可能只值 0.001 美元，也就是 0.1 分钱。但走传统信用卡要扣 0.30 美元，30 分钱手续费。**付款比买的东西贵了 300 倍**——这个市场根本上不了规模。

所有按次计费的 Pay-per-use 市场——比如 API 调用、AI 推理调用、按字数计费的内容、按次的数据查询——都被传统支付的单位经济学**彻底锁死**几十年。Agent 经济需要一个新支付层。

而 x402 + 稳定币 USDC on Base L2 解锁了这个市场：

L2 结算速度约 2 秒，Flashblock 预确认 200 毫秒就能给反馈，传统 ACH/Card 要 T+2-3 天。
单笔成本不到 1 分钱，让小于 1 美元的微支付变得理所当然。
零账户设置——有钱包就能付，不用注册商户。
无 API Key 风险——签名是一次性授权，泄露也无害。
天然全球同价——不管在哪个国家，费率都一样。

让小于 1 美元的 Agent 微支付**从不可能变成理所当然**。这就是 AWS、Stripe、Cloudflare 同时押注这个协议的真正原因——不是因为它酷，而是因为它解锁了一个被锁住几十年的万亿级市场。""")
ftr(s, 7)

# ============================================================================
#  SAVE
# ============================================================================
out = Path(__file__).parent / "openagentpay-talk-p1.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
print(f"   Total slides: {len(prs.slides)}")

# Verify notes
total_chars = 0
for i, slide in enumerate(prs.slides, 1):
    notes = slide.notes_slide.notes_text_frame.text
    chars = len(notes)
    total_chars += chars
    print(f"   Slide {i}: notes = {chars} chars")
print(f"   Total notes: {total_chars} chars")
