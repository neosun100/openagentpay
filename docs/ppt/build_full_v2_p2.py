"""
build_full_v2_p2.py — 接续 build_full_v2_p1，添加 slide 11-23

读取 openagentpay-talk-half.pptx (10 slides)
追加 13 张 slide → 最终 23 页 → 保存为 openagentpay-talk.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ===== 配色 + helper (与 p1 相同) =====
BG       = RGBColor(0x0A, 0x1A, 0x2E)
BG_CARD  = RGBColor(0x14, 0x2A, 0x44)
BG_DARK  = RGBColor(0x05, 0x10, 0x1F)
FG       = RGBColor(0xF0, 0xF4, 0xF8)
DIM      = RGBColor(0x8B, 0xA0, 0xBA)
MUTED    = RGBColor(0x4A, 0x5C, 0x75)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)
GOLD     = RGBColor(0xFF, 0xB8, 0x00)
AWS_ORG  = RGBColor(0xFF, 0x99, 0x00)
OAP_GRN  = RGBColor(0x00, 0xFF, 0x88)
RED      = RGBColor(0xFF, 0x44, 0x44)
PURPLE   = RGBColor(0xBF, 0x5A, 0xF2)

IMG_DIR = Path(__file__).parent / "images"
TOTAL_PAGES = 23

def bg(slide, color=BG):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color
def t(slide, left, top, w, h, text, *, size=18, color=FG, bold=False, align="left", font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb
def mt(slide, left, top, w, h, lines, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb
def bar(slide, left, top, w, h, color=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s
def card(slide, left, top, w, h, color=BG_CARD, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h); s.adjustments[0] = 0.08
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s
def hdr(slide, title, subtitle=None, accent=CYAN):
    t(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6), title, size=26, color=FG, bold=True)
    if subtitle:
        t(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4), subtitle, size=13, color=DIM)
    bar(slide, Inches(0.6), Inches(1.4), Inches(2), Emu(40000), accent)
def ftr(slide, page_num=None):
    t(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
      "AgentCore Payments + OpenAgentPay  ·  2026-05-18  ·  Neo Sun", size=9, color=MUTED)
    if page_num:
        t(slide, Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.3),
          f"{page_num:02d} / {TOTAL_PAGES:02d}", size=9, color=MUTED, align="right")
def n(slide, text):
    slide.notes_slide.notes_text_frame.text = text
def metric(slide, left, top, w, h, number, label, color=CYAN, num_size=44, lbl_size=11):
    card(slide, left, top, w, h, BG_CARD, color)
    t(slide, left, top + Inches(0.3), w, Inches(1.2), number, size=num_size, color=color, bold=True, align="center")
    t(slide, left, top + h - Inches(0.5), w, Inches(0.3), label, size=lbl_size, color=DIM, align="center")
def chap(prs, num, title, subtitle, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG_DARK)
    t(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6), f"第 {num} 章", size=24, color=accent, bold=True)
    t(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.5), title, size=52, color=FG, bold=True)
    t(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.6), subtitle, size=18, color=DIM)
    bar(s, Inches(0.6), Inches(5.4), Inches(3), Emu(60000), accent)
    return s

# Load p1 result
prs = Presentation(str(Path(__file__).parent / "openagentpay-talk-half.pptx"))
print(f"Loaded {len(prs.slides)} slides from half")

# ===== Slide 11: 两大阵营 + AWS 战略选择 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "AWS 为什么只做买方侧  ·  整篇文章最战略级洞察",
    "买方侧 (AWS+Coinbase+Stripe) vs 收款侧 (Visa+MC+PayPal+Google)", AWS_ORG)

img_path = str(IMG_DIR / "04-strategic-positioning.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(7.5), Inches(4.2))

# 右半：4 个理由
card(s, Inches(8.4), Inches(1.7), Inches(4.4), Inches(4.2), BG_CARD, AWS_ORG)
t(s, Inches(8.6), Inches(1.85), Inches(4.0), Inches(0.4),
  "AWS 4 个战略理由", size=13, color=AWS_ORG, bold=True)
mt(s, Inches(8.6), Inches(2.3), Inches(4.0), Inches(3.4), [
    ("1. 避开卡网络强势区", 11, FG, True),
    ("Visa/MC 卡网络几十年", 10, DIM, False),
    ("硬卷肯定输", 10, DIM, False),
    ("", 6, DIM, False),
    ("2. 紧贴 AWS 算力优势", 11, FG, True),
    ("买方侧每分价值变 Bedrock", 10, DIM, False),
    ("/存储/流量收入", 10, DIM, False),
    ("", 6, DIM, False),
    ("3. 双下注", 11, FG, True),
    ("Coinbase=Crypto, Stripe=Fiat", 10, DIM, False),
    ("不赌谁赢", 10, DIM, False),
    ("", 6, DIM, False),
    ("4. 开发者心智", 11, FG, True),
    ("Stripe+Coinbase 是开发者", 10, DIM, False),
    ("最信任的支付品牌", 10, DIM, False),
])

# 底部金句
card(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.85), BG_CARD, GOLD)
t(s, Inches(0.85), Inches(6.15), Inches(11.7), Inches(0.4),
  "Visa/Mastercard/PayPal 在建「商户接受」基础设施", size=12, color=DIM, bold=False)
t(s, Inches(0.85), Inches(6.5), Inches(11.7), Inches(0.4),
  "AWS 在建「Agent 发起」基础设施  —  两边互补，不是替代", size=14, color=GOLD, bold=True)

n(s, """【两大阵营 + AWS 战略】

兄弟们这是整篇文章最战略级的洞察——所有人都在卷 Agent 支付，**但站位不同**。

看左边图——**两大阵营**：

**左边（Agent 买方侧）**：AWS + Coinbase + Stripe — "让 Agent 能方便地**去买东西**"

**右边（Merchant 收款侧）**：Visa + Mastercard + PayPal + Google — "让商户能**接受** Agent 发起的交易"

这是**两个完全不同的战场**。

**AWS 为什么选左边？4 个战略理由**：

1. **避开卡网络强势区**——Visa / Mastercard 经营卡网络几十年，硬卷肯定输
2. **紧贴 AWS 自身优势**——开发者基础 + 算力平台，买方侧的每一分价值都变成 Bedrock / 存储 / 流量的收入
3. **双下注**——Coinbase 押 Crypto-native，Stripe Privy 押 Fiat + Card，**不用赌谁赢**
4. **开发者心智**——Stripe + Coinbase 都是开发者**信任度最高**的支付品牌（相对于传统银行）

**底部金句很关键**——The AI Economy 的 Ken Yeung 一针见血：

> "Visa, Mastercard, and PayPal are building infrastructure for merchants to **ACCEPT** agent payments. AWS's approach is to create infrastructure for AI agents to **MAKE** those purchases."

（Visa/MC/PayPal 在建"商户接受"基础设施；AWS 在建"Agent 发起"基础设施。）

**两边互补，不是替代**。

**对我们 SA 的意义**：跟客户讲这个产品时不要把它定位成"跟 Visa 抢生意"——是"补 Visa 没做的那一面"。卡网络 + AWS 一起把 Agent 经济跑起来。""")
ftr(s, 11)

# ===== Slide 12: 7 层 Payment Guardrail =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "7 层 Payment Guardrail  ·  企业级安全设计",
    "Agent 不会乱花钱  ·  基础设施层硬强制  ·  Agent 无法绕过", PURPLE)

img_path = str(IMG_DIR / "06-security-layers.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(7.5), Inches(5.0))

# 右半：7 层列表
card(s, Inches(8.4), Inches(1.7), Inches(4.5), Inches(5.0), BG_CARD, PURPLE)
t(s, Inches(8.6), Inches(1.85), Inches(4.1), Inches(0.4),
  "7 层硬强制", size=13, color=PURPLE, bold=True)
mt(s, Inches(8.6), Inches(2.3), Inches(4.1), Inches(4.4), [
    ("1. Authorization Layer", 10, FG, True),
    ("终端用户必须显式授权", 9, DIM, False),
    ("", 5, DIM, False),
    ("2. Session Layer", 10, FG, True),
    ("maxSpendAmount + expiryTime", 9, DIM, False),
    ("", 5, DIM, False),
    ("3. Policy Layer", 10, FG, True),
    ("per-agent / per-session 策略", 9, DIM, False),
    ("", 5, DIM, False),
    ("4. On-chain Layer", 10, FG, True),
    ("不可篡改的链上记录", 9, DIM, False),
    ("", 5, DIM, False),
    ("5. Compliance Layer", 10, FG, True),
    ("Sanctions + 反洗钱", 9, DIM, False),
    ("", 5, DIM, False),
    ("6. Identity Layer", 10, FG, True),
    ("Agent 拿不到 private key", 9, DIM, False),
    ("", 5, DIM, False),
    ("7. Audit Layer", 10, FG, True),
    ("SOX/MRM 监管审计", 9, DIM, False),
])

# 底部强调
card(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4), BG_CARD, PURPLE)
t(s, Inches(0.8), Inches(6.9), Inches(12), Inches(0.3),
  '"Agent never has open-ended access to funds. It operates only within defined limits."  ←  AWS 官方原话',
  size=11, color=PURPLE, bold=True, align="center")

n(s, """【7 层 Payment Guardrail】

兄弟们这页特别重要——**这是每个 CFO/CTO 听到"Agent 自动付款"的第一反应**：

> "Agent 被 prompt injection 骗了怎么办？一晚上把钱花光怎么办？审计怎么做？"

AWS 的答案是 **7 层 Payment Guardrail**——在基础设施层硬强制，**Agent 无法绕过**。

**每一笔支付必须过 7 道关**（看图，从外到内）：

1. **Authorization Layer**——终端用户必须**显式授权**，Agent 不能自主绑钱包
2. **Session Layer**——maxSpendAmount + expiryTime 硬限额，达到即拒绝
3. **Policy Layer**——财务/合规团队可定义 per-agent / per-session 精细策略
4. **On-chain Layer**——每笔产生**不可篡改**的链上记录（区块链特性）
5. **Compliance Layer**——Coinbase CDP Facilitator 内置 sanctions + 反洗钱检查
6. **Identity Layer**——AgentCore Identity 隔离密钥，**Agent 拿不到 private key**
7. **Audit Layer**——支付决策 + 推理链 + tx hash 统一日志，支持 SOX / MRM / 金融监管审计

**任一层失败 → 支付立即阻断**。这就是为什么 Warner Bros Discovery、Heurist AI 这些公司敢接入的底气。

**AWS 官方原话**："Agent never has open-ended access to funds. It operates only with explicit permission and within defined limits."

**对我们 SA 的意义**：客户问"Agent 不安全吧"——答案就是这 7 层。这是**企业可接受 Agent 付款的关键设计**——金融客户尤其看重。Heurist 之所以敢上生产，因为有这 7 层兜底。

兄弟们记住这一页——客户问"Agent 安全吗"，先讲 Identity Layer（Agent 拿不到私钥）+ Session Layer（硬限额），再讲剩下 5 层。""")
ftr(s, 12)

# ===== Slide 13: 三阶段路线图 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "AgentCore Payments 三阶段路线图",
    "Preview → GA → Expansion  ·  开发者体验跨阶段一致", AWS_ORG)

img_path = str(IMG_DIR / "07-roadmap.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(8.5), Inches(4.5))

# 右半：阶段说明
card(s, Inches(9.3), Inches(1.7), Inches(3.6), Inches(1.45), BG_CARD, CYAN)
t(s, Inches(9.5), Inches(1.85), Inches(3.3), Inches(0.4),
  "Phase 1 · Preview", size=12, color=CYAN, bold=True)
t(s, Inches(9.5), Inches(2.25), Inches(3.3), Inches(0.4),
  "今天 (2026-05)", size=10, color=DIM)
mt(s, Inches(9.5), Inches(2.55), Inches(3.3), Inches(0.55), [
    ("· 微支付 < $1/call", 9, FG, False),
    ("· 4 region · x402 only", 9, FG, False),
])

card(s, Inches(9.3), Inches(3.25), Inches(3.6), Inches(1.45), BG_CARD, GOLD)
t(s, Inches(9.5), Inches(3.4), Inches(3.3), Inches(0.4),
  "Phase 2 · GA", size=12, color=GOLD, bold=True)
t(s, Inches(9.5), Inches(3.8), Inches(3.3), Inches(0.4),
  "下一步", size=10, color=DIM)
mt(s, Inches(9.5), Inches(4.1), Inches(3.3), Inches(0.55), [
    ("· 加 Fiat (Stripe 全球)", 9, FG, False),
    ("· 更多 region · MPP 协议", 9, FG, False),
])

card(s, Inches(9.3), Inches(4.8), Inches(3.6), Inches(1.45), BG_CARD, OAP_GRN)
t(s, Inches(9.5), Inches(4.95), Inches(3.3), Inches(0.4),
  "Phase 3 · Expansion", size=12, color=OAP_GRN, bold=True)
t(s, Inches(9.5), Inches(5.35), Inches(3.3), Inches(0.4),
  "未来", size=10, color=DIM)
mt(s, Inches(9.5), Inches(5.65), Inches(3.3), Inches(0.55), [
    ("· Agent 订机票/酒店", 9, FG, False),
    ("· 跨商户购物", 9, FG, False),
])

# 底部金句
card(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5), BG_CARD, GOLD)
t(s, Inches(0.85), Inches(6.45), Inches(11.7), Inches(0.4),
  '"开发者体验跨阶段一致。今天写的代码到阶段 3 不用重写"  ←  AWS 承诺',
  size=12, color=GOLD, bold=True, align="center")

n(s, """【三阶段路线图】

AWS 已经明示了产品演进路径——这 3 个阶段。

**Phase 1 · Preview（今天）**：
- 微支付 < $1/call
- 4 个 Region (us-east-1/us-west-2/eu-west-1/ap-southeast-2)
- 只有 x402 协议
- 支持 Coinbase CDP + Stripe Privy 两个钱包

**Phase 2 · GA**：
- 加 Fiat 法币支付（Stripe 全球轨道）
- 更多 Region（亚洲会进来吗？路线图没明确说）
- 可能加 MPP 协议
- 这是 Stripe 真正的舞台——传统支付场景

**Phase 3 · Expansion**：
- **Agent 帮你订机票、预订酒店、跨商户购物** —— Agent-on-behalf-of-buyer
- 这是真正的"Agent 经济"——Agent 替你做完整的购物决策

**AWS 关键承诺**（金色那行）：

> "The developer experience stays consistent across each phase. Configure your wallet, set your policies, and your agent transacts."

翻译：**你今天写的代码，到阶段 3 不用重写**。AWS 承担协议升级、新钱包接入的复杂度。

**对我们 SA 的意义**：
- 客户问"现在能用吗"——只能用 Preview 微支付场景，订机票还得等
- 客户问"什么时候 GA"——AWS 没给时间表，但金融服务客户应该是首批
- 客户问"未来能干嘛"——Phase 3 的"跨商户购物"才是真正的 game changer，那时候 Agent 就是经济主体

**对 OpenAgentPay 的意义**：我们的扩展是为 Phase 1/2 的客户解决"亚洲 region 缺失"和"钱包覆盖不足"的问题。等 AWS 自己做到 Phase 3，我们的 connector 也是 day-1 reference impl。""")
ftr(s, 13)

# ===== Slide 14: 客户案例 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "已经有哪些客户在用",
    "Heurist 生产 ✓  ·  Warner Bros 评估中  ·  Cox/Thomson Reuters/PGA 即将升级", AWS_ORG)

# 4 张客户卡片
# Heurist - 生产
card(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(2.5), BG_CARD, OAP_GRN)
t(s, Inches(0.85), Inches(1.95), Inches(5.5), Inches(0.4), "Heurist AI", size=20, color=OAP_GRN, bold=True)
t(s, Inches(0.85), Inches(2.35), Inches(5.5), Inches(0.4), "已生产上线 ✓  ·  金融研究 Agent", size=12, color=DIM)
mt(s, Inches(0.85), Inches(2.75), Inches(5.5), Inches(1.5), [
    ("\"用 AgentCore Payments 给金融研究 Agent 加付费", 11, FG, False),
    ("能力。终端客户设预算，Agent 自动购买实时行情、", 11, FG, False),
    ("社交舆情、新闻数据。", 11, FG, False),
    ("", 6, DIM, False),
    ("几行代码就接入了。\"  — JW Wang, 创始人", 11, GOLD, True),
])

# Warner Bros - 评估中
card(s, Inches(6.7), Inches(1.8), Inches(6.0), Inches(2.5), BG_CARD, GOLD)
t(s, Inches(6.95), Inches(1.95), Inches(5.5), Inches(0.4), "Warner Bros Discovery", size=20, color=GOLD, bold=True)
t(s, Inches(6.95), Inches(2.35), Inches(5.5), Inches(0.4), "积极评估中  ·  顶级内容即时付费消费", size=12, color=DIM)
mt(s, Inches(6.95), Inches(2.75), Inches(5.5), Inches(1.5), [
    ("\"在积极探索 AgentCore Payments —— 想让顶级", 11, FG, False),
    ("内容（体育赛事、大片）能在 Agent 互动中即时", 11, FG, False),
    ("付费消费。\"", 11, FG, False),
    ("", 6, DIM, False),
    ("— Mit Majithia, 执行副总裁", 11, GOLD, True),
])

# AgentCore 既有客户群（即将升级）
card(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.4), BG_CARD, CYAN)
t(s, Inches(0.85), Inches(4.55), Inches(11.5), Inches(0.4),
  "AgentCore 平台现有客户  ·  潜在升级", size=18, color=CYAN, bold=True)
t(s, Inches(0.85), Inches(4.95), Inches(11.5), Inches(0.4),
  '已用 AgentCore 平台做复杂 workflow  ·  "those agents can also transact"', size=12, color=DIM)

# 三个 logo-style 卡片
card(s, Inches(0.85), Inches(5.55), Inches(3.9), Inches(1.1), BG_DARK, MUTED)
t(s, Inches(0.85), Inches(5.7), Inches(3.9), Inches(0.4), "Cox Automotive", size=14, color=FG, bold=True, align="center")
t(s, Inches(0.85), Inches(6.05), Inches(3.9), Inches(0.4), "汽车销售平台", size=11, color=DIM, align="center")

card(s, Inches(4.85), Inches(5.55), Inches(3.9), Inches(1.1), BG_DARK, MUTED)
t(s, Inches(4.85), Inches(5.7), Inches(3.9), Inches(0.4), "Thomson Reuters", size=14, color=FG, bold=True, align="center")
t(s, Inches(4.85), Inches(6.05), Inches(3.9), Inches(0.4), "金融数据 + 法律研究", size=11, color=DIM, align="center")

card(s, Inches(8.85), Inches(5.55), Inches(3.85), Inches(1.1), BG_DARK, MUTED)
t(s, Inches(8.85), Inches(5.7), Inches(3.85), Inches(0.4), "PGA TOUR", size=14, color=FG, bold=True, align="center")
t(s, Inches(8.85), Inches(6.05), Inches(3.85), Inches(0.4), "体育赛事内容", size=11, color=DIM, align="center")

n(s, """【客户案例】

AWS 官方 Blog 里明确点名的客户分三类：

**🏆 已生产上线（Heurist AI）**
JW Wang 创始人原话："Heurist 已经在用 AgentCore Payments 做金融和加密研究 Agent。终端客户给研究设预算，Agent 自动购买实时行情、社交舆情、新闻数据。**几行代码就接入了**。"

这是**最有说服力的 PR**——客户现身说法 + production 验证。

**🎬 积极评估中（Warner Bros Discovery）**
Mit Majithia 执行副总裁："在积极探索 AgentCore Payments——想让顶级内容（体育赛事、大片）能在 Agent 互动中即时付费消费。"

WBD 的兴趣点是**内容订阅 + 即时付费**——比如你看体育直播突然弹出付费高清回放，Agent 替你按需付。

**🔄 AgentCore 平台现有客户（潜在升级）**
- Cox Automotive（汽车销售平台 Agent）
- Thomson Reuters（金融数据 + 法律研究 Agent）
- PGA TOUR（体育赛事内容 Agent）

这三家**已经在用 AgentCore 平台**做复杂 workflow。今天的新能力意味着："those agents can also transact"——这些 Agent 现在也能交易了。换言之，它们是**最容易被引导升级的第一批客户**。

**对我们 SA 的意义**：
- Heurist 是金融场景客户案例，可以直接给我们的金融客户讲
- WBD 是内容场景，可以给我们的媒体客户讲
- Cox/Thomson Reuters/PGA 是 AgentCore 平台升级路径，我们 SA 主要是**找已经用 AgentCore 平台的客户**升级到带 Payments 版本

**Heurist 案例特别值钱**——它证明**几行代码就能集成**这件事是真的，不是 PR 话术。""")
ftr(s, 14)

print(f"After slides 11-14: {len(prs.slides)} slides")

# ===== Slide 15: 章节 2 分隔（短板与转折）=====
s = chap(prs, 2, "短板与转折", "好产品，但当前覆盖不到我们的客户  ·  这是机会", RED)
n(s, """【第二章 · 短板与转折】

刚才讲完 AgentCore Payments 11 页——非常完整的产品深度解读。兄弟们应该认同：**这是个 promising 的好产品**。

但任何 Preview 产品都有覆盖面有限的问题，AWS 也不可能一上来就支持所有场景。接下来 2 页，我从我们 Web3 团队的角度看：哪些客户能用、哪些不能用、机会在哪里。

**短板不是 AWS 产品的过错，是我们 Web3 团队的机会**。""")
ftr(s, 15)

# ===== Slide 16: 短板 (图 09) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "AgentCore Payments 现状  ·  好产品，但 4 个限制",
    "Coinbase + Stripe only  ·  4 region 无亚洲  ·  BYO 接口未开放", RED)

img_path = str(IMG_DIR / "09-wallet-comparison.png")
s.shapes.add_picture(img_path, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.0))

# 4 个限制
mt(s, Inches(0.6), Inches(1.8), Inches(6), Inches(0.5), [
    ("当前 4 个限制（按重要性）", 14, RED, True),
])

card(s, Inches(0.6), Inches(2.3), Inches(6.2), Inches(1.0), BG_CARD, RED)
t(s, Inches(0.85), Inches(2.4), Inches(5.7), Inches(0.4), "1. 仅 2 个钱包", size=14, color=FG, bold=True)
t(s, Inches(0.85), Inches(2.8), Inches(5.7), Inches(0.4), "Coinbase CDP + Stripe Privy  ·  无 HashKey/Binance/OKX", size=11, color=DIM)

card(s, Inches(0.6), Inches(3.4), Inches(6.2), Inches(1.0), BG_CARD, RED)
t(s, Inches(0.85), Inches(3.5), Inches(5.7), Inches(0.4), "2. 仅 1 个协议", size=14, color=FG, bold=True)
t(s, Inches(0.85), Inches(3.9), Inches(5.7), Inches(0.4), "x402 v1/v2 only  ·  没接 MPP/AP2/ACP", size=11, color=DIM)

card(s, Inches(0.6), Inches(4.5), Inches(6.2), Inches(1.0), BG_CARD, RED)
t(s, Inches(0.85), Inches(4.6), Inches(5.7), Inches(0.4), "3. 仅 4 region 可用", size=14, color=FG, bold=True)
t(s, Inches(0.85), Inches(5.0), Inches(5.7), Inches(0.4), "美/欧/澳  ·  注意：没有亚洲！", size=11, color=RED, bold=True)

card(s, Inches(0.6), Inches(5.6), Inches(6.2), Inches(1.0), BG_CARD, RED)
t(s, Inches(0.85), Inches(5.7), Inches(5.7), Inches(0.4), "4. BYO connector 接口未开放", size=14, color=FG, bold=True)
t(s, Inches(0.85), Inches(6.1), Inches(5.7), Inches(0.4), "第三方钱包当前没法 register 进去", size=11, color=DIM)

n(s, """【现状与短板】

AgentCore Payments 当前 4 个限制（按对我们 Web3 团队的重要性排序）：

**1. 仅 2 个钱包**——Coinbase CDP + Stripe Privy
看似覆盖很广（一个 custodial、一个 self-custody），但仔细看：
- Coinbase 是美国上市公司，主要服务北美 Web3 用户
- Stripe Privy 是 2024 Stripe 收购的，主要服务北美和欧洲合规支付场景
- 我们覆盖的 HashKey、Binance Pay、OKX、Bitget、Bybit、HashKey Pro——一个都不在列表里

**2. 仅 1 个协议**——x402 only
没接 MPP（Stripe + Tempo）、没接 AP2（Google）、没接 ACP（OpenAI）。客户用其他协议时只能用别的方案。

**3. 仅 4 个 region 可用**——这是最致命的短板！
- us-east-1, us-west-2, eu-west-1, ap-southeast-2
- **没有亚洲 region**！新加坡、东京、香港、孟买，AgentCore Payments 都用不了
- ap-southeast-2 是悉尼，对中国大陆客户延迟 200ms+

**4. BYO connector 接口未开放**——Preview 阶段限制
意味着**第三方钱包没法 register 进去**——HashKey 不能自己写个 connector 接进去。

兄弟们到这里应该有感觉——这个产品**在亚洲场景、CEX 场景，当前都不能用**。

下一页讲我们的客户具体是哪些——这就是 OpenAgentPay 的起点。""")
ftr(s, 16)

# ===== Slide 17: 转折页 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_DARK)
t(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
  "那个关键问题", size=14, color=DIM)
t(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.5),
  "我们的客户在哪里？", size=64, color=GOLD, bold=True)
t(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.4),
  "AWS Web3 Team 服务的客户群体  (none of these works on AgentCore Payments today)", size=14, color=DIM)

card(s, Inches(0.6),  Inches(3.7), Inches(2.95), Inches(2.2), BG_CARD, RED)
t(s, Inches(0.8), Inches(3.85), Inches(2.7), Inches(0.4), "亚洲合规交易所", size=13, color=RED, bold=True)
t(s, Inches(0.8), Inches(4.3), Inches(2.7), Inches(1.5),
  "HashKey · OKX\nBitget · Bybit\nGate.io", size=14, color=FG)

card(s, Inches(3.75), Inches(3.7), Inches(2.95), Inches(2.2), BG_CARD, RED)
t(s, Inches(3.95), Inches(3.85), Inches(2.7), Inches(0.4), "亚洲 CEX Pay", size=13, color=RED, bold=True)
t(s, Inches(3.95), Inches(4.3), Inches(2.7), Inches(1.5),
  "Binance Pay\nOKX Pay\nBitget Wallet", size=14, color=FG)

card(s, Inches(6.9), Inches(3.7), Inches(2.95), Inches(2.2), BG_CARD, RED)
t(s, Inches(7.1), Inches(3.85), Inches(2.7), Inches(0.4), "Web3 自托管钱包", size=13, color=RED, bold=True)
t(s, Inches(7.1), Inches(4.3), Inches(2.7), Inches(1.5),
  "MetaMask\nWalletConnect\nRainbow", size=14, color=FG)

card(s, Inches(10.05), Inches(3.7), Inches(2.65), Inches(2.2), BG_CARD, RED)
t(s, Inches(10.25), Inches(3.85), Inches(2.4), Inches(0.4), "传统亚洲支付", size=13, color=RED, bold=True)
t(s, Inches(10.25), Inches(4.3), Inches(2.4), Inches(1.5),
  "支付宝\n微信支付\nUnionPay", size=14, color=FG)

card(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7), BG_DARK, GOLD)
t(s, Inches(0.8), Inches(6.18), Inches(12), Inches(0.55),
  "全部用不上 AgentCore Payments  —  这不是产品的错，是我们要解决的问题。",
  size=15, color=GOLD, bold=True, align="center")

n(s, """【转折页】

所以问题来了：我们 AWS Web3 团队，服务的客户群体——亚洲合规交易所、亚洲 CEX Pay、Web3 自托管钱包、传统亚洲支付——**4 个类别全部当前用不上 AgentCore Payments**。

这不是 AWS 产品做得不好。Coinbase + Stripe 已经覆盖了北美和欧洲的主要场景。这是 AWS 推产品的策略选择——先做最大公分母，等市场反馈再扩张。

但这就把我们 Web3 团队卡死了。

我们去找 HashKey 推 AgentCore Payments，HashKey 会问："你们支持 HashKey Chain 吗？支持 HKDR 港币稳定币吗？"答：暂时不支持。
找 Binance 推，会问："支持 Binance Pay API 吗？"答：不支持。
OKX、Bitget、Bybit、支付宝、微信...同样的问题。

兄弟们，我们怎么办？

**3 个选项**：

**选项 1：等 AWS 自己开放支持**——但 roadmap 写的是 'Others* coming soon'，没时间表，我们等不起客户。

**选项 2：放弃**——跟客户说 AWS 这个东西好但你用不上。但作为 AWS SA，这个 narrative 太弱了。

**选项 3：自己做扩展层**——让我们的客户也能用。**这是我选的路径——OpenAgentPay**。

但**单兵作战做不成**——所以今天来邀请兄弟们一起做。接下来 5 页讲清楚 OpenAgentPay 是什么、当前到哪里、未来兄弟们能怎么参与。""")
ftr(s, 17)

print(f"After slides 15-17: {len(prs.slides)} slides")

# ===== Slide 18: 章节 3 分隔（OpenAgentPay）=====
s = chap(prs, 3, "OpenAgentPay 扩展", "AgentCore Payments 的开放扩展层  ·  Work-in-Progress  ·  抛砖引玉", OAP_GRN)
n(s, """【第三章 · OpenAgentPay 扩展】

接下来 5 页讲 OpenAgentPay——我做的扩展项目。

**特别注意一个心态**——这是 work-in-progress，**不是 production，不是 done deal**。我会诚实地告诉哪里完成了、哪里还在开发、哪里需要兄弟们加入。

但当前的进度已经有一些**有真东西的成果**：协议层验证完成、链上 e2e 跑通、AWS 部署 live、4 笔真上链 tx 在 Blockscout 永久可查。

**我说"邀请兄弟们参与"不是空话**——后面有具体的参与方式。

兄弟们听这部分可以**比第一部分更放松**——这部分是讨论 + 邀请，不是我宣讲。""")
ftr(s, 18)

# ===== Slide 19: 5 层架构 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "OpenAgentPay  ·  5 层可插拔架构",
    "类比 Kubernetes CRI/CSI/CNI  ·  Agent Payments 的可插拔时刻", OAP_GRN)
img_path = str(IMG_DIR / "platform-architecture.png")
s.shapes.add_picture(img_path, Inches(2.3), Inches(1.7), Inches(8.7), Inches(5.0))
card(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4), BG_CARD, GOLD)
t(s, Inches(0.8), Inches(6.9), Inches(12), Inches(0.3),
  "核心承诺：业务代码改 1 行  →  payment_manager = PaymentManager(wallet_provider=\"hashkey-chain\")",
  size=12, color=GOLD, bold=True, align="center")
n(s, """【OpenAgentPay 5 层架构】

OpenAgentPay 一句话定位：**AgentCore Payments 的开放可插拔扩展层**——让任何钱包、任何协议、任何稳定币都能即插即用接入。

我用一个**餐厅类比**让兄弟们一秒理解 5 层架构：

- Layer 1 · Strands Plugin → **服务员的点单 PAD**（Agent 直接用）
- Layer 2 · Payment Orchestrator → **餐厅经理**（编排订单 + 检查预算 + 协议路由）
- Layer 3 · Protocols → **支付方式**（信用卡 / 微信 / 现金 / 链上, 可选）
- Layer 4 · Wallet Connectors → **POS 机**（每种支付方式一台，按统一接口）
- Layer 5 · Self-Hosted Facilitator → **银行后台清算系统**

**5 层的核心设计原则**：Layer 2/3/5 是 framework 不变；Layer 1 是 SDK 入口不变；**只有 Layer 4 按钱包变化**。

**用具体例子说明这个设计的价值**：

假如今天 HashKey 工程师想给 Agent 加付款能力：
- ❌ 不用 OpenAgentPay：自己实现协议+钱包+上链+错误处理+安全 — **1-2 个月 + 1 工程师**
- ✅ 用 OpenAgentPay：在 Layer 4 写 HashKeyChainConnector — **1-2 天 + 1 工程师**

新加 OKX？同样在 Layer 4 写一个 OKXConnector，**1-2 天**。其他层不动。

**类比 Kubernetes**：CRI（Container Runtime Interface）/ CSI / CNI 让任何容器、存储、网络可插拔接入。**OpenAgentPay 想做 Agent Payments 的 CRI 时刻**。

底部金色字是**核心承诺**：业务代码改 1 行 → 切换钱包：
- payment_manager = PaymentManager(wallet_provider="hashkey-chain")
- payment_manager = PaymentManager(wallet_provider="binance-pay")
- payment_manager = PaymentManager(wallet_provider="coinbase-cdp")

**兄弟们如果想参与的话——主要是 Layer 4**。每人选一个钱包做 connector，1-2 天 ship 一个 demo。这就是 community effort 的具体形态。""")
ftr(s, 19)

# ===== Slide 20: 双协议轨道 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "双协议轨道  ·  为什么不能一律用 x402",
    "x402 是链上协议  ·  CEX 结构上不上链  ·  协议形状共享，加密层可插拔", OAP_GRN)
img_path = str(IMG_DIR / "protocol-comparison.png")
s.shapes.add_picture(img_path, Inches(2.3), Inches(1.7), Inches(8.7), Inches(5.0))
card(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4), BG_CARD, GOLD)
t(s, Inches(0.8), Inches(6.9), Inches(12), Inches(0.3),
  "强行套 x402 等于让中心化数据库假装成 ERC-20 合约  ·  既笨拙又违背 CEX 的低成本优势",
  size=12, color=GOLD, bold=True, align="center")
n(s, """【双协议轨道 · OAP-CEX 命名 + 含义】

这页讲一个项目最容易被 challenge 的核心问题：**为什么 Binance 不直接用 x402？**

**先讲 OAP-CEX 名字怎么来的**：

**OAP-CEX = Open**Agent**P**ay - **C**entralized **EX**change
- OAP 是项目品牌缩写
- CEX 是用途（Centralized Exchange）
- 借鉴 x.509 / x402 / TLS 等工业界传统——**短前缀 + 用途后缀**
- 未来扩展：OAP-FIAT（传统支付）、OAP-BANK（银行直连）

**回到核心问题：为什么 Binance 不能用 x402？**

简短回答：**x402 是链上协议，但 Binance 这类 CEX 结构上不上链**。

兄弟们想想——你在 Binance 看到的"1000 USDT 余额"，**不是链上 USDT 合约的状态，是 Binance 内部数据库里的一行记录**：user_id=12345, asset=USDT, balance=1000。

CEX 的真相：**99% 的资金流动在内部账本，1% 在链上（提现时才上链）**。

- 内部转账：Binance 数据库 update，~50ms 完成，**0 gas**
- 提现到链：才上链，要 gas，几秒确认

**为什么 CEX 这么设计**？成本（链上转账要 gas）+ 速度（链确认几秒 vs 毫秒）+ 合规（KYC 信息只能在 CEX 内部）+ 效率（高频用户全上链链就堵了）。

**强行套 x402 给 Binance 等于让中心化数据库假装成 ERC-20 合约——既笨拙又违背 CEX 的低成本优势**。

**OpenAgentPay 的关键 insight**：x402 协议的"形状"很好（402 challenge → sign → retry），但"加密层 + 结算层"应该可插拔。所以我们：

- **x402** = 协议形状 + EIP-712 签名层（链上钱包用）
- **OAP-CEX** = 同样协议形状 + HMAC 签名层（CEX 用）

**为什么客户会选 OAP-CEX？**
- HashKey、Binance、OKX 客户：他们的用户**已经在 CEX**，不需要折腾链
- 传统支付场景（支付宝、微信、Stripe credit card）：根本不上链
- B2B 高额场景：KYC 必需，CEX 路径更合规

**杀招总结**："x402 让 Web3 user 能付款，OAP-CEX 让剩下 90% 的用户也能付款"。""")
ftr(s, 20)

# ===== Slide 21: Live AWS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "Live on AWS  ·  完整合规架构",
    "Browser → CloudFront → API Gateway → Lambda → HashKey Chain", OAP_GRN)
img_path = str(IMG_DIR / "architecture.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(7.0), Inches(5.0))
card(s, Inches(8.0), Inches(1.7), Inches(4.7), Inches(1.5), BG_CARD, OAP_GRN)
t(s, Inches(8.2), Inches(1.85), Inches(4.4), Inches(0.4), "Live Demo URL", size=14, color=OAP_GRN, bold=True)
t(s, Inches(8.2), Inches(2.3), Inches(4.4), Inches(0.4), "d1p7yxa99nxaye", size=24, color=FG, bold=True)
t(s, Inches(8.2), Inches(2.7), Inches(4.4), Inches(0.4), ".cloudfront.net", size=24, color=FG, bold=True)
metric(s, Inches(8.0), Inches(3.4), Inches(2.3), Inches(1.2), "1s",     "cold start", CYAN, num_size=28)
metric(s, Inches(10.4), Inches(3.4), Inches(2.3), Inches(1.2), "5s",   "end-to-end", GOLD, num_size=28)
metric(s, Inches(8.0), Inches(4.7), Inches(2.3), Inches(1.2), "357s",   "CDK deploy", PURPLE, num_size=28)
metric(s, Inches(10.4), Inches(4.7), Inches(2.3), Inches(1.2), "82",   "tests pass", OAP_GRN, num_size=28)
card(s, Inches(8.0), Inches(6.0), Inches(4.7), Inches(0.7), BG_CARD, OAP_GRN)
t(s, Inches(8.2), Inches(6.05), Inches(4.4), Inches(0.3), "Compliance", size=11, color=OAP_GRN, bold=True)
t(s, Inches(8.2), Inches(6.3), Inches(4.4), Inches(0.4), "no Function URL · IAM-scoped · KMS at rest", size=10, color=DIM)

n(s, """【Live on AWS · 完整合规架构】

OpenAgentPay 不是 ppt-only project——demo 已经部署到 AWS 生产环境。

【完整链路】
1. 用户浏览器 → HTTPS POST /api/pay
2. CloudFront（CDN + DDoS + TLS + CORS）→ /api/* 路由到 API Gateway
3. API Gateway HTTP API（合规公网入口）→ 转 Lambda integration
4. Lambda（Node 20）→ 查 Secrets Manager 拿私钥 → viem 构造 EIP-712 签名
5. Lambda → HashKey Chain RPC broadcast → 等 ~5 秒确认
6. 返回 tx hash → API Gateway → CloudFront → 用户浏览器显示 ✅

【4 个合规原则】
1. No Function URL（Palisade-proof）
2. IAM-scoped throughout
3. KMS at rest
4. Secrets never in logs

══════════════════════════════════════════
【🎬 DEMO 演讲台操作流程】
══════════════════════════════════════════

**这一页停留时切到浏览器：https://d1p7yxa99nxaye.cloudfront.net**

Step 1：在 Tab 1 'Run Demo' 停留，给兄弟们看 4 步流程概览
Step 2：点 Step 1 'Run' → 显示链上 USDC 余额（约 5987 USDC）
       说："这都是从 HashKey Chain 真链上读的"
Step 3：点 Step 2 'Run' → 创建 Payment Session（spend governor 边界）
Step 4：点 Step 3 'Pay' → EIP-712 签名 + 上链结算，等 5 秒
       说："看，tx 0x...拿到了"
Step 5：点 Blockscout 链接 → 浏览器跳转到链上证据
Step 6（如果时间够）：切 Tab 3 'AI Agent' → 点付费按钮真上链

**演讲前 30 分钟跑这两条命令拿现场数据**：
```bash
# HSK gas 余额
curl -s -X POST https://testnet.hsk.xyz -H "Content-Type: application/json" \\
  -d '{"jsonrpc":"2.0","method":"eth_getBalance","params":["0x863d9C87b6bBd4aEf115C297C41643a0b887eAd7","latest"],"id":1}' \\
  | python3 -c "import sys,json; print(f'HSK: {int(json.load(sys.stdin)[\\"result\\"], 16) / 10**18:.10f}')"

# USDC 业务余额
curl -s -X POST https://testnet.hsk.xyz -H "Content-Type: application/json" \\
  -d '{"jsonrpc":"2.0","method":"eth_call","params":[{"to":"0x0685C487Df4Cc0723Aa828C299686798294E9803","data":"0x70a08231000000000000000000000000863d9c87b6bbd4aef115c297c41643a0b887ead7"},"latest"],"id":1}' \\
  | python3 -c "import sys,json; print(f'USDC: {int(json.load(sys.stdin)[\\"result\\"], 16) / 10**6:.6f}')"
```""")
ftr(s, 21)

# ===== Slide 22: HashKey demo（HSK + USDC 真相）=====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "HashKey Chain 链上 demo  ·  HSK 和 USDC 真相",
    "5 步可复现  ·  4 笔真 tx 链上可查  ·  双语言独立实现", OAP_GRN)

# 左半：5 步 + 关键数据
mt(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.5), [
    ("第 1 步  ·  faucet 领 0.1 HSK (gas 用)", 13, CYAN, True),
])
t(s, Inches(0.8), Inches(2.05), Inches(5.7), Inches(0.3),
  "https://faucet.hashkeychain.net/faucet", size=10, color=DIM, font="Consolas")

mt(s, Inches(0.6), Inches(2.45), Inches(6.0), Inches(0.5), [
    ("第 2 步  ·  部署 MockUSDC 合约 (EIP-3009)", 13, CYAN, True),
])
t(s, Inches(0.8), Inches(2.8), Inches(5.7), Inches(0.3),
  "0x0685C487Df4Cc0723Aa828C299686798294E9803", size=10, color=DIM, font="Consolas")

mt(s, Inches(0.6), Inches(3.2), Inches(6.0), Inches(0.5), [
    ("第 3 步  ·  Mint USDC 到 agent 钱包", 13, CYAN, True),
])

mt(s, Inches(0.6), Inches(3.65), Inches(6.0), Inches(0.5), [
    ("第 4 步  ·  EIP-712 签名 + transferWithAuthorization", 13, CYAN, True),
])

mt(s, Inches(0.6), Inches(4.1), Inches(6.0), Inches(0.5), [
    ("第 5 步  ·  Python ref + TypeScript prod 双实现", 13, CYAN, True),
])

# HSK + USDC 关系卡
card(s, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.0), BG_CARD, GOLD)
t(s, Inches(0.85), Inches(4.85), Inches(5.5), Inches(0.4),
  "HSK + USDC 关系（必懂）", size=12, color=GOLD, bold=True)
mt(s, Inches(0.85), Inches(5.25), Inches(5.5), Inches(1.4), [
    ("HSK = 原生币 (付 gas，类比 ETH)", 11, FG, False),
    ("USDC = 我们部署的 ERC-20 (业务结算)", 11, FG, False),
    ("", 6, DIM, False),
    ("等 Circle 部署官方 USDC，换合约地址即生产级", 10, DIM, False),
    ("协议层 0 改动 → production-equivalent", 10, OAP_GRN, True),
])

# 右半：4 笔链上 tx 卡片
card(s, Inches(7.0), Inches(1.7), Inches(5.7), Inches(5.0), BG_CARD, OAP_GRN)
t(s, Inches(7.25), Inches(1.85), Inches(5.2), Inches(0.4),
  "✅ 链上事实证据 (Blockscout 永久可查)", size=13, color=OAP_GRN, bold=True)
mt(s, Inches(7.25), Inches(2.35), Inches(5.2), Inches(4.3), [
    ("合约部署 tx", 11, CYAN, True),
    ("0xb9bdfdb1a975413dab1825824a88...", 9, DIM, False),
    ("", 6, DIM, False),
    ("Python e2e tx", 11, CYAN, True),
    ("0xff8a175e3f4b41a30b67940a4b65...", 9, DIM, False),
    ("", 6, DIM, False),
    ("TypeScript e2e tx", 11, CYAN, True),
    ("0x5c10e2ae5a152169c5870ce440f7...", 9, DIM, False),
    ("", 6, DIM, False),
    ("AWS Lambda 生产 tx", 11, OAP_GRN, True),
    ("0x4562d26ea86ad19f303f378bcfbf...", 9, DIM, False),
    ("", 8, DIM, False),
    ("两套独立实现产生相同链上效果", 11, GOLD, True),
    ("→ 协议层抽象正确的最强证据", 10, DIM, False),
])

n(s, """【HashKey demo · HSK 和 USDC 真相】

【HSK 和 USDC 关系（讲清楚）】

很多兄弟会问：从 faucet 领的是 HSK，但 demo 显示 USDC，这是怎么回事？

**HSK 和 USDC 是两个完全不同的代币，扮演两个完全不同的角色，跟以太坊上 ETH 和 USDT 关系一模一样**：

类比以太坊：ETH (gas) ←→ USDT/USDC (业务转账)
类比 HashKey Chain：HSK (gas) ←→ MockUSDC (业务结算)

**5 步还原**：
1. 从 faucet 领 0.1 HSK（链原生币，付 gas）
2. 用 HSK 当 gas 部署 MockUSDC 合约（ERC-20 + EIP-3009）
3. 调 mint() 给自己钱包 mint 1000+ USDC
4. EIP-712 签名 transferWithAuthorization
5. Python 参考实现 + TypeScript 生产实现，独立两套

**关键 insight**：USDC 是真上链的、签名真做、合约真跑、gas 真烧——只有 mint 这部分是我们自己 mint 的（mock token 嘛）。**等 Circle 在 HashKey Chain 部官方 USDC，换合约地址 = 生产级**。

**右边 4 笔链上 tx**：
- 合约部署 tx
- Python e2e tx ✅
- TypeScript e2e tx ✅（同样链上效果，证明协议层抽象正确）
- AWS Lambda 生产 tx ✅（最新生产环境）

**演讲台杀招话术**：
> "如果质疑这是 mock——我刚查了链上：HSK 从 0.1 减到 0.0999977，烧了 0.00000207 HSK 总 gas，16 笔交易全部上链。每笔 gas 0.0000001 HSK——这个数字太低没法手动 mint，必须真上链才有这个 cost。"

**给兄弟们的会后验证**：testnet-explorer.hsk.xyz 输合约地址 0x0685C487 → 看到所有 tx → 链上 immutable 永远可查。""")
ftr(s, 22)

# ===== Slide 23: 三种集成形态 + CTA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "三种集成形态 + 邀请兄弟们一起做",
    "Standalone → Plugin → AgentCore Native  ·  让所有客户都能用 Payments", OAP_GRN)

# 三种形态卡片
card(s, Inches(0.6), Inches(1.7), Inches(4.0), Inches(2.4), BG_CARD, CYAN)
t(s, Inches(0.85), Inches(1.85), Inches(3.7), Inches(0.4), "形态 1", size=12, color=CYAN, bold=True)
t(s, Inches(0.85), Inches(2.25), Inches(3.7), Inches(0.5), "Standalone", size=20, color=FG, bold=True)
t(s, Inches(0.85), Inches(2.7), Inches(3.7), Inches(0.4), "✅ 当前 live", size=11, color=OAP_GRN, bold=True)
mt(s, Inches(0.85), Inches(3.1), Inches(3.7), Inches(1.0), [
    ("CloudFront + Lambda", 11, DIM, False),
    ("不依赖 Bedrock", 11, DIM, False),
    ("适用：Web3 创业 / HashKey", 10, DIM, False),
])

card(s, Inches(4.7), Inches(1.7), Inches(4.0), Inches(2.4), BG_CARD, GOLD)
t(s, Inches(4.95), Inches(1.85), Inches(3.7), Inches(0.4), "形态 2", size=12, color=GOLD, bold=True)
t(s, Inches(4.95), Inches(2.25), Inches(3.7), Inches(0.5), "Strands Plugin", size=20, color=FG, bold=True)
t(s, Inches(4.95), Inches(2.7), Inches(3.7), Inches(0.4), "🚧 下一步 1-2 周", size=11, color=GOLD, bold=True)
mt(s, Inches(4.95), Inches(3.1), Inches(3.7), Inches(1.0), [
    ("AgentCore Runtime", 11, DIM, False),
    ("Identity / Memory / Obs.", 11, DIM, False),
    ("适用：Bedrock + Strands 客户", 10, DIM, False),
])

card(s, Inches(8.8), Inches(1.7), Inches(4.0), Inches(2.4), BG_CARD, OAP_GRN)
t(s, Inches(9.05), Inches(1.85), Inches(3.7), Inches(0.4), "形态 3", size=12, color=OAP_GRN, bold=True)
t(s, Inches(9.05), Inches(2.25), Inches(3.7), Inches(0.5), "AgentCore Native", size=20, color=FG, bold=True)
t(s, Inches(9.05), Inches(2.7), Inches(3.7), Inches(0.4), "⏳ 等 AWS 开放 BYO", size=11, color=PURPLE, bold=True)
mt(s, Inches(9.05), Inches(3.1), Inches(3.7), Inches(1.0), [
    ("BYO connector", 11, DIM, False),
    ("AWS Console 内置", 11, DIM, False),
    ("适用：全部 AWS 客户", 10, DIM, False),
])

# CTA 大块
card(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.5), BG_CARD, OAP_GRN)
t(s, Inches(0.85), Inches(4.45), Inches(11.5), Inches(0.5),
  "邀请兄弟们一起做 · 让所有 AWS 亚太客户都能用 Payments",
  size=20, color=OAP_GRN, bold=True, align="center")

mt(s, Inches(0.85), Inches(5.0), Inches(11.5), Inches(1.7), [
    ("· 谁的客户有 Payments 需求 → 找我聊，1-2 天 ship 一个 connector", 12, FG, False),
    ("· 谁负责覆盖某个交易所/钱包 → 一起接 OKX / Bitget / Bybit / HashKey Pro", 12, FG, False),
    ("· 用了发现 bug 或建议 → GitHub 提 issue / PR", 12, FG, False),
    ("· 一起对外发声 → FSI DNB SA Team 整体 community contribution", 12, FG, False),
    ("", 8, DIM, False),
    ("github.com/neosun100/openAgentPay  ·  Apache 2.0  ·  jiasunm@amazon.com", 12, GOLD, True),
])

# 底部一句
t(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
  "个人力量很渺小  ·  团队拧成一股绳，每人 1-2 天，2 周覆盖整个亚太市场  ·  Q&A 时间欢迎打断",
  size=11, color=DIM, align="center", bold=True)

n(s, """【三种集成形态 + CTA · 邀请兄弟们一起做】

**先快速过三种集成形态**——让任何客户都能 self-identify 自己是哪种。

**形态 1：Standalone**（当前 live）
独立 Lambda + CloudFront 部署。客户即使不用 Bedrock 也能用。
适用：Web3 创业公司、HashKey 自家 Agent 平台。

**形态 2：Strands Plugin**（下一步 1-2 周）
作为 Strands Plugin 接入 AgentCore Runtime + Identity + Memory + Observability。业务代码 1 行不改：把 AgentCorePaymentsPlugin 换成 OpenAgentPayPlugin。
适用：用 Bedrock + Strands 的客户（最大客户群）。

**形态 3：AgentCore Native**（等 AWS 开放 BYO connector）
当 AWS 开放 BYO 接口，我们直接 register。客户在 AgentCore Console 里下拉选 HashKey Chain、Binance Pay、OKX Pay。
适用：全部 AWS 客户（主流市场）。

══════════════════════════════════════════
【真正的 CTA · 兄弟们一起做】
══════════════════════════════════════════

**说实话——一个人接不完所有钱包**。

我列一下当前需要兄弟们参与的具体方式：

**第一种：负责一个钱包/协议方向**
- 覆盖 OKX 客户的兄弟 → 一起做 OKXConnector，1-2 天
- 覆盖 Bitget / Bybit / HashKey Pro 的兄弟 → 同上
- Binance 已经我做了 work-in-progress

**第二种：客户带场景进来**
你的客户问起 Payments，找我聊：
- 客户用 MetaMask → 1 天写个 MetaMaskConnector
- 客户用 Stripe → 1 天 + 1 天加传统支付协议轨道
- 客户在做 RWA → 把 OpenAgentPay 当成支付层 demo 给客户

**第三种：架构 review + 反馈**
OpenAgentPay 的协议层、framework、安全设计——兄弟们看了有想法直接提 issue / PR。开源 Apache 2.0。

**第四种：一起对外发声**
把这个工作输出成 **FSI DNB SA Team 的 community contribution**——内部分享、外部博客、再到 AWS Solutions Library。这就是我们团队的 visibility。

══════════════════════════════════════════
【最后一句话】
══════════════════════════════════════════

兄弟们，**个人力量很渺小**——一个人接一个 connector 1-2 天，但接 20 个钱包就 1-2 个月。**团队拧成一股绳，每人 1-2 天，2 周就能覆盖整个亚太市场**。

这就是为什么我说"邀请兄弟们一起做"——不是客套，是实打实的需求。

**终极意义**：让所有 AWS 客户都能用 Agent Payments，不只是北美 Coinbase + Stripe 用户。这是我们 AWS Web3 SA 团队该做的事，作为 FSI DNB SA Team 的整体贡献。

**谢谢兄弟们**。Q&A 时间——任何问题、想法、批评、challenge 都欢迎。

GitHub: github.com/neosun100/openAgentPay
Live demo: d1p7yxa99nxaye.cloudfront.net
私聊: jiasunm@amazon.com""")
ftr(s, 23)

# ============================================================================
#  SAVE
# ============================================================================
out = Path(__file__).parent / "openagentpay-talk.pptx"
prs.save(out)
print(f"\n✅ Saved final: {out}")
print(f"   Total slides: {len(prs.slides)}")

# Verify all notes
total_chars = 0
for i, slide in enumerate(prs.slides, 1):
    notes = slide.notes_slide.notes_text_frame.text
    chars = len(notes)
    total_chars += chars
    title_text = "(no title)"
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            title_text = shape.text_frame.text.strip().split("\n")[0][:50]
            break
    print(f"   Slide {i:2d}  notes={chars:4d}  {title_text}")
print(f"\n   📊 Total notes: {total_chars} chars  ({total_chars / len(prs.slides):.0f} avg/slide)")
