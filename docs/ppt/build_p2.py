"""
build_p2.py — OpenAgentPay × AgentCore Payments 周会分享 PPT (Slide 8-16)

读 build_p1 输出，append slides 8-16:
  8. 章节分隔（第 2 章：短板与转折）
  9. 现状 + 短板（图 09）
  10. 转折页（"我们的客户在哪里？"）
  11. 章节分隔（第 3 章：OpenAgentPay 解决方案）
  12. OpenAgentPay 5 层架构（platform-architecture.png）
  13. 双协议轨道（protocol-comparison.png）
  14. Live AWS（architecture.png + Live URL + demo 演示指引）
  15. HashKey 链上 demo（faucet 截图 + 4 笔真 tx）
  16. 三种集成形态 + CTA（合并）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ============================================================================
#  CONFIG (与 p1 一致)
# ============================================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG       = RGBColor(0x0A, 0x1A, 0x2E)
BG_CARD  = RGBColor(0x14, 0x2A, 0x44)
BG_DARK  = RGBColor(0x05, 0x10, 0x1F)
FG       = RGBColor(0xF0, 0xF4, 0xF8)
DIM      = RGBColor(0x8B, 0xA0, 0xBA)
MUTED    = RGBColor(0x4A, 0x5C, 0x75)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)
GOLD     = RGBColor(0xFF, 0xB8, 0x00)
AWS_ORG  = RGBColor(0xFF, 0x99, 0x00)
OAP_GRN  = RGBColor(0x00, 0xFF, 0x88)
RED      = RGBColor(0xFF, 0x44, 0x44)
PURPLE   = RGBColor(0xBF, 0x5A, 0xF2)

IMG_DIR = Path(__file__).parent / "images"

# ============================================================================
#  HELPER FUNCTIONS (与 p1 一致，内联)
# ============================================================================
def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def t(slide, left, top, w, h, text, *, size=18, color=FG, bold=False, align="left", font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb

def mt(slide, left, top, w, h, lines, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb

def bar(slide, left, top, w, h, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def card(slide, left, top, w, h, color=BG_CARD, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def hdr(slide, title, subtitle=None, accent=CYAN):
    t(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
      title, size=28, color=FG, bold=True)
    if subtitle:
        t(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
          subtitle, size=14, color=DIM)
    bar(slide, Inches(0.6), Inches(1.45), Inches(2), Emu(40000), accent)

def ftr(slide, page_num=None, total=16):
    t(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
      "OpenAgentPay × AgentCore Payments  ·  2026-05-18  ·  Neo Sun",
      size=9, color=MUTED)
    if page_num:
        t(slide, Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.3),
          f"{page_num:02d} / {total:02d}", size=9, color=MUTED, align="right")

def n(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text

def metric(slide, left, top, w, h, number, label, color=CYAN, num_size=44, lbl_size=11):
    card(slide, left, top, w, h, BG_CARD, color)
    t(slide, left, top + Inches(0.3), w, Inches(1.2),
      number, size=num_size, color=color, bold=True, align="center")
    t(slide, left, top + h - Inches(0.5), w, Inches(0.3),
      label, size=lbl_size, color=DIM, align="center")

def chap(prs, num, title, subtitle, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG_DARK)
    t(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
      f"第 {num} 章", size=24, color=accent, bold=True)
    t(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.5),
      title, size=56, color=FG, bold=True)
    t(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.6),
      subtitle, size=18, color=DIM)
    bar(s, Inches(0.6), Inches(5.5), Inches(3), Emu(60000), accent)
    return s

# ============================================================================
#  LOAD p1, append p2 slides
# ============================================================================
prs = Presentation(str(Path(__file__).parent / "openagentpay-talk-p1.pptx"))

# Update p1 slides ftr 总页数 (从 14 改成 16)
# Skip — 已经在 p1 里写死了，新页码用 16 即可

# ============================================================================
# Slide 8 — 章节分隔：第 2 章 短板与转折
# ============================================================================
s = chap(prs, 2, "现状与短板", "Preview 阶段 · 只支持 2 个钱包 · 客户在哪里？", RED)
n(s, """【第二章】上一章我们看到 AgentCore Payments 是个非常 promising 的产品——技术对、协议对、生态对、客户案例真实。但它当前 Preview 阶段有一个**关键限制**，对我们 Web3 团队来说是致命的。

接下来 2 页，我把数据摆出来——不是为了 trash 这个产品，而是为了看清楚我们的位置：哪些客户能用、哪些不能用、机会在哪里。""")
ftr(s, 8)

# ============================================================================
# Slide 9 — 现状 + 短板（图 09）
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "AgentCore Payments 现状  ·  好产品，但只支持 2 个钱包",
    "x402 生态成熟  ·  AgentCore Payments 当前只接通了 Coinbase + Stripe", RED)

# 左半：4 个 metric 数据
t(s, Inches(0.6), Inches(1.8), Inches(6), Inches(0.4),
  "x402 生态  (协议层非常成熟)", size=14, color=CYAN, bold=True)

metric(s, Inches(0.6),  Inches(2.3), Inches(2.85), Inches(1.4), "75M",   "近 30 天 tx", CYAN, num_size=28)
metric(s, Inches(3.7),  Inches(2.3), Inches(2.85), Inches(1.4), "590K+", "活跃买家",     OAP_GRN, num_size=28)

metric(s, Inches(0.6),  Inches(3.85), Inches(2.85), Inches(1.4), "100K+", "活跃卖家",     PURPLE, num_size=28)
metric(s, Inches(3.7),  Inches(3.85), Inches(2.85), Inches(1.4), "<$0.01", "单笔成本",   GOLD, num_size=28)

# 短板列表
t(s, Inches(0.6), Inches(5.5), Inches(6), Inches(0.4),
  "AgentCore Payments 限制  (产品层短板)", size=14, color=RED, bold=True)
mt(s, Inches(0.6), Inches(5.95), Inches(6), Inches(1.2), [
    ("•  仅 2 个钱包：Coinbase CDP · Stripe Privy", 12, FG, False),
    ("•  仅 1 个协议：x402 v1/v2  (没接 MPP/AP2/ACP)", 12, FG, False),
    ("•  仅 4 个 region (无亚洲！)", 12, RED, True),
    ("•  Preview 阶段，BYO connector 接口未开放", 12, FG, False),
])

# 右半：图 09 wallet comparison
img_path = str(IMG_DIR / "09-wallet-comparison.png")
s.shapes.add_picture(img_path, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.0))

n(s, """【现状与短板】这页把数据摆出来。

左边是 x402 生态的现状——非常成熟。30 天 7541 万笔交易、活跃买家 59 万、卖家 10 万、单笔成本不到 1 分钱。这是一个生产级、被头部玩家信赖的协议层。

右边这张图是 AgentCore Payments 当前支持的两种钱包路径：

第一种 Coinbase CDP——Custodial Wallet 模式，钱包私钥由 Coinbase 托管。优势是 fast onboarding，缺点是用户依赖 Coinbase。
第二种 Stripe Privy——Embedded Wallet 模式，每个用户一个 EOA 钱包，私钥分片在用户设备 + Privy 后端。

看似覆盖很广——一个 custodial、一个 self-custody。但仔细看：

**Coinbase 是美国上市公司**，主要服务北美 Web3 用户。
**Stripe Privy 是 2024 年 Stripe 收购的**，主要服务北美和欧洲合规支付场景。

下面这 4 条限制特别关键：
1. 仅 2 个钱包；
2. 仅 1 个协议（x402），AgentCore 还没接 MPP/AP2 等其他正在涌现的 Agent 支付协议；
3. **仅 4 个 region 可用**（us-east-1, us-west-2, eu-west-1, ap-southeast-2）——注意，**没有亚洲 region**。新加坡、东京、香港、孟买，AgentCore Payments 都用不了；
4. Preview 阶段，BYO connector 接口未开放——意味着第三方钱包没法 register 进去。

这就是问题。下一页讲我们的客户在哪里。""")
ftr(s, 9)

# ============================================================================
# Slide 10 — 转折页：我们的客户在哪？
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_DARK)

# 大问题
t(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
  "那个关键问题", size=14, color=DIM)
t(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.5),
  "我们的客户在哪里？", size=64, color=GOLD, bold=True)

# 4 个客户类别
t(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.4),
  "AWS Web3 Team 服务的客户群体  (none of these works on AgentCore Payments today)", size=14, color=DIM)

card(s, Inches(0.6),  Inches(3.7), Inches(2.95), Inches(2.2), BG_CARD, RED)
t(s, Inches(0.8), Inches(3.85), Inches(2.7), Inches(0.4), "亚洲合规交易所", size=13, color=RED, bold=True)
t(s, Inches(0.8), Inches(4.3), Inches(2.7), Inches(1.5),
  "HashKey · OKX\nBitget · Bybit\nGate.io", size=14, color=FG)

card(s, Inches(3.75), Inches(3.7), Inches(2.95), Inches(2.2), BG_CARD, RED)
t(s, Inches(3.95), Inches(3.85), Inches(2.7), Inches(0.4), "亚洲 CEX Pay", size=13, color=RED, bold=True)
t(s, Inches(3.95), Inches(4.3), Inches(2.7), Inches(1.5),
  "Binance Pay\nOKX Pay\nBitget Wallet", size=14, color=FG)

card(s, Inches(6.9), Inches(3.7), Inches(2.95), Inches(2.2), BG_CARD, RED)
t(s, Inches(7.1), Inches(3.85), Inches(2.7), Inches(0.4), "Web3 自托管钱包", size=13, color=RED, bold=True)
t(s, Inches(7.1), Inches(4.3), Inches(2.7), Inches(1.5),
  "MetaMask\nWalletConnect\nRainbow", size=14, color=FG)

card(s, Inches(10.05), Inches(3.7), Inches(2.65), Inches(2.2), BG_CARD, RED)
t(s, Inches(10.25), Inches(3.85), Inches(2.4), Inches(0.4), "传统亚洲支付", size=13, color=RED, bold=True)
t(s, Inches(10.25), Inches(4.3), Inches(2.4), Inches(1.5),
  "支付宝\n微信支付\nUnionPay", size=14, color=FG)

# 底部金句
card(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7), BG_DARK, GOLD)
t(s, Inches(0.8), Inches(6.18), Inches(12), Inches(0.55),
  "全部用不上 AgentCore Payments  —  这不是产品的错，是我们要解决的问题。",
  size=15, color=GOLD, bold=True, align="center")

n(s, """【转折页】所以问题来了：我们 AWS Web3 团队，服务的客户群体——亚洲合规交易所、亚洲 CEX Pay、Web3 自托管钱包、传统亚洲支付——全部 4 个类别，**当前都用不上 AgentCore Payments**。

这不是 AWS 产品做得不好。Coinbase + Stripe 已经覆盖了北美和欧洲的主要场景。这是 AWS 推产品的策略选择——先做最大公分母，等市场反馈再扩张。

但这就把我们 Web3 团队卡死了。我们去找 HashKey 推 AgentCore Payments，HashKey 会问：'你们支持 HashKey Chain 吗？'答：不支持。
我们去找 Binance 推，Binance 会问：'你们支持 Binance Pay API 吗？'答：不支持。
我们去找 OKX、Bitget、Bybit……同样的问题。

我们 Web3 团队的核心 KPI 是覆盖亚太 Web3 客户。但手里这个产品，**目标客户一个都用不上**。

我们能怎么办？

选项 1：等 AWS 自己开放支持。但 roadmap 写的是 'Others* coming soon' ——什么时候、是否包含 HashKey、AWS 没给时间表，**我们等不起**。
选项 2：放弃这个产品，跟客户说 AWS 这个东西好但你用不上。这个 narrative 太弱了，AWS SA 不能这么做。
选项 3：自己做扩展层，让我们的客户也能用。**这是我选的路径——OpenAgentPay**。

接下来我用 6 页讲清楚 OpenAgentPay 是什么、怎么做、跑得通吗、未来怎么走。""")
ftr(s, 10)

# ============================================================================
# Slide 11 — 章节分隔：第 3 章 OpenAgentPay 解决方案
# ============================================================================
s = chap(prs, 3, "OpenAgentPay 解决方案",
         "AgentCore Payments 的开放扩展层  ·  让任何钱包、任何协议、任何稳定币都能即插即用",
         OAP_GRN)
n(s, """【第三章】接下来 5 页讲我做的扩展方案——OpenAgentPay。

请大家注意一个心态——这是 work-in-progress，不是 production。我会诚实地告诉大家哪里完成了、哪里还在开发、哪里需要大家加入。

但当前的进度已经很有说服力了：协议层验证完成、链上 demo 跑通、AWS 部署 live、4 笔真上链 tx 可查。这是个**有真东西的 work-in-progress**。""")
ftr(s, 11)

# ============================================================================
# Slide 12 — OpenAgentPay 5 层架构
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "OpenAgentPay  ·  5 层可插拔架构",
    "类比 Kubernetes 之于容器  ·  CRI/CSI/CNI 标准化  →  Agent Payments 的 CRI 时刻", OAP_GRN)

# 嵌入 platform-architecture 图（占大半屏）
img_path = str(IMG_DIR / "platform-architecture.png")
s.shapes.add_picture(img_path, Inches(2.3), Inches(1.7), Inches(8.7), Inches(5.0))

# 底部承诺
card(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4), BG_CARD, GOLD)
t(s, Inches(0.8), Inches(6.9), Inches(12), Inches(0.3),
  "核心承诺：业务代码改 1 行  →  payment_manager = PaymentManager(wallet_provider=\"hashkey-chain\")",
  size=12, color=GOLD, bold=True, align="center")

n(s, """【OpenAgentPay 5 层架构】OpenAgentPay 一句话定位：AgentCore Payments 的开放可插拔扩展层，让任何钱包、任何协议、任何稳定币都能即插即用接入。

这张图是 5 层架构，从上到下：

**Layer 1：Strands Plugin** — Python SDK，drop-in 替换 AgentCore 原版的 AgentCorePaymentsPlugin。客户业务代码 0 改动，只换 plugin 配置。

**Layer 2：Payment Orchestrator** — PaymentManager / SessionManager (spend governor + TTL) / ConnectorRegistry / ProtocolRouter。命名跟 AgentCore Payments 1:1 对齐——这是关键，意味着我们跟 AWS 不竞争，是扩展。

**Layer 3：Protocols（可插拔）** — 当前已支持 x402 v1/v2（链上）和 OAP-CEX v0.1（CEX，我自己设计的）。未来 MPP / AP2 / ACP 都可加。

**Layer 4：Wallet Connectors（可插拔）** — 当前 HashKey Chain（live）+ Binance Pay 已实现。下方灰色字列出未来要接的：MetaMask、WalletConnect、OKX、Bitget、Bybit、HashKey Pro。按需触发，不写时间。

**Layer 5：Self-Hosted Facilitator** — CloudFront → API Gateway → Lambda → KMS 私钥管理。

类比 Kubernetes：CRI (Container Runtime Interface) / CSI / CNI 让任何容器、存储、网络可插拔接入。OpenAgentPay 想做 **Agent Payments 的 CRI 时刻**——定义 WalletConnector + ProtocolAdapter 标准接口。

底部那行金色字是核心承诺：业务代码改 1 行，从 hashkey-chain 切到 binance-pay 切到 coinbase-cdp。框架的真正价值就在这一行代码。""")
ftr(s, 12)

# ============================================================================
# Slide 13 — 双协议轨道
# ============================================================================
s = prs.slides.add_slide(prs.slide_invert if False else prs.slide_layouts[6])
bg(s, BG)
hdr(s, "双协议轨道  ·  为什么不能一律用 x402",
    "x402 是链上协议，但 Binance 这类 CEX 结构上不上链  ·  协议形状共享，加密层可插拔", OAP_GRN)

# 嵌入 protocol-comparison 图（核心图）
img_path = str(IMG_DIR / "protocol-comparison.png")
s.shapes.add_picture(img_path, Inches(2.3), Inches(1.7), Inches(8.7), Inches(5.0))

# 底部金句
card(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4), BG_CARD, GOLD)
t(s, Inches(0.8), Inches(6.9), Inches(12), Inches(0.3),
  "强行套 x402 等于让中心化数据库假装成 ERC-20 合约  ·  既笨拙又违背 CEX 的低成本优势",
  size=12, color=GOLD, bold=True, align="center")

n(s, """【双协议轨道】这页讲一个深刻的问题——这是项目最容易被问到的核心 architectural question：**为什么 Binance 不直接用 x402？**

简短回答：x402 是**链上协议**，但 Binance 这类 CEX **结构上不上链**。强行套 x402 等于让中心化数据库假装成 ERC-20 合约——既笨拙又违背 CEX 的低成本优势。

看左侧 x402 链上路径：EIP-712 签名 + ECDSA、EIP-3009 链上结算、0x... 以太坊地址、Facilitator 替 Agent 付 gas、约 5 秒 L2 finalize。
看右侧 OAP-CEX 路径：HMAC-SHA256/512 签名（CEX 标准）、CEX 内部账本结算 off-chain、内部 merchant ID、不上链所以没 gas、约 50 毫秒结算。

两个**世界观完全不同**：链上世界讲 self-custody + permissionless + on-chain truth；CEX 世界讲 KYC + 中心化账本 + 低延迟。

OpenAgentPay 的**关键 insight**：x402 协议的'形状'很好（402 challenge → sign → retry），但'加密层 + 结算层'**应该可插拔**。所以我们拆分：

x402 = 协议形状 + EIP-712 签名层（链上钱包用）
OAP-CEX = 同样协议形状 + HMAC 签名层（CEX 用）

两者**共享 ProtocolAdapter 接口**；PaymentManager 通过 ProtocolRouter 自动按 402 response signature 派发。

这就是为什么我们既能跑通 HashKey Chain 链上结算（x402），也能跑通 Binance Pay CEX 内部结算（OAP-CEX）——**用同一套 PaymentManager**。

OAP-CEX 是我自己写的协议规范（packages/protocol-cex-pay/doc/SPEC.md，24 页 IETF-style draft，向后兼容 x402）。下一步是推到 IETF/W3C 标准化。

这就是 OpenAgentPay 跟 AgentCore 的真正差异化——AgentCore 只走链上路径，我们双轨。""")
ftr(s, 13)

# ============================================================================
# Slide 14 — Live AWS 部署架构
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "Live on AWS  ·  完整合规架构",
    "Browser → CloudFront → API Gateway → Lambda → HashKey Chain  ·  全链路真实运行", OAP_GRN)

# 嵌入 architecture 图
img_path = str(IMG_DIR / "architecture.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(7.0), Inches(5.0))

# 右侧：Live URL + 数据
card(s, Inches(8.0), Inches(1.7), Inches(4.7), Inches(1.5), BG_CARD, OAP_GRN)
t(s, Inches(8.2), Inches(1.85), Inches(4.4), Inches(0.4),
  "Live Demo URL", size=14, color=OAP_GRN, bold=True)
t(s, Inches(8.2), Inches(2.3), Inches(4.4), Inches(0.4),
  "d1p7yxa99nxaye", size=24, color=FG, bold=True)
t(s, Inches(8.2), Inches(2.7), Inches(4.4), Inches(0.4),
  ".cloudfront.net", size=24, color=FG, bold=True)

# Performance metrics
metric(s, Inches(8.0), Inches(3.4), Inches(2.3), Inches(1.2), "1s",     "cold start", CYAN, num_size=28)
metric(s, Inches(10.4), Inches(3.4), Inches(2.3), Inches(1.2), "5s",   "end-to-end", GOLD, num_size=28)

metric(s, Inches(8.0), Inches(4.7), Inches(2.3), Inches(1.2), "357s",   "CDK deploy", PURPLE, num_size=28)
metric(s, Inches(10.4), Inches(4.7), Inches(2.3), Inches(1.2), "82",   "tests pass", OAP_GRN, num_size=28)

# Compliance highlights
card(s, Inches(8.0), Inches(6.0), Inches(4.7), Inches(0.7), BG_CARD, OAP_GRN)
t(s, Inches(8.2), Inches(6.05), Inches(4.4), Inches(0.3),
  "Compliance", size=11, color=OAP_GRN, bold=True)
t(s, Inches(8.2), Inches(6.3), Inches(4.4), Inches(0.4),
  "no Function URL · IAM-scoped · KMS at rest", size=10, color=DIM)

n(s, """【Live on AWS】OpenAgentPay 不是 ppt-only project——我们今晚刚把 demo 部署到 AWS 生产环境。

看这张架构图：Browser → CloudFront (CDN + DDoS) → API Gateway HTTP API → Lambda (Node 20) → HashKey Chain Testnet。侧边依赖：Secrets Manager（KMS 加密私钥）+ HashKey Chain RPC。

这里有一个**今晚特别重要的 ops 故事**：我们最初用 Lambda Function URL，被 Amazon Palisade 自动检测为 'world accessible'，Epoxy 自动 mitigation 把权限收窄。后来改成 API Gateway HTTP API——这是 AWS 标准合规公网入口，不会被 flag。这个 incident 反而展现了 AWS secure-by-default 的能力。

右边这 4 个 metric：1 秒冷启动、5 秒端到端结算、357 秒 CDK 一键部署、82 个 tests 全过。Compliance：no Function URL（Palisade-proof）、IAM-scoped throughout、KMS at rest。

【DEMO 演示指引 - 演讲台操作流程】

此时打开浏览器到 https://d1p7yxa99nxaye.cloudfront.net

**演示步骤（建议 3-4 分钟）**：

1. 先在 Tab 1 'Run Demo' 停留——给大家看 4 步流程概览
2. 点 Step 1 'Run' → 实时显示 HashKey Chain 链上 USDC 余额（约 998 USDC）
   说："这都是从 HashKey Chain 真链上读的，不是 mock"
3. 点 Step 2 'Run' → 创建 Payment Session（预算 $1，TTL 60 min）
   说："这是 spend governor 边界，超预算硬拒绝，基础设施层强制"
4. 点 Step 3 'Pay' → 真上链结算 0.001 USDC
   说："现在 EIP-712 typed data 签名 + Facilitator 上链，等 5 秒"
5. 看到 tx hash → 点击 Blockscout 链接 → 浏览器跳转
   说："这就是真上链的证据，链上 immutable 永远可查"
6. 切到 Tab 3 'AI Agent'（如果时间够）→ 点付费按钮（"ETH 深度分析" 或 "减半研报"）
   说："这是模拟 Strands Agent 自主决策——免费工具直接用，付费工具触发链上结算"

如果 demo 卡了：备好 https://github.com/neosun100/openAgentPay 截图 + Blockscout 历史 tx 链接 fallback。""")
ftr(s, 14)

# ============================================================================
# Slide 15 — HashKey 链上 demo 流程
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "HashKey Chain 链上 demo  ·  完整 5 步可复现",
    "从 faucet 领测试币 → 部署合约 → mint → 签名 → 上链  ·  4 笔真 tx 已在 Blockscout 可查", OAP_GRN)

# 左半：5 步流程
mt(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.5), [
    ("第 1 步  ·  生成 EVM 钱包  (viem)", 13, CYAN, True),
])
t(s, Inches(0.8), Inches(2.05), Inches(5.7), Inches(0.3),
  "Agent 地址 0x863d9C87b6bBd4aEf115C297C41643a0b887eAd7", size=10, color=DIM, font="Consolas")

mt(s, Inches(0.6), Inches(2.45), Inches(6.0), Inches(0.5), [
    ("第 2 步  ·  HashKey Testnet faucet 领 HSK gas", 13, CYAN, True),
])
t(s, Inches(0.8), Inches(2.8), Inches(5.7), Inches(0.3),
  "https://faucet.hashkeychain.net/faucet  →  约 1 HSK", size=10, color=DIM, font="Consolas")

mt(s, Inches(0.6), Inches(3.2), Inches(6.0), Inches(0.5), [
    ("第 3 步  ·  部署 MockUSDC 合约 (EIP-3009 完整实现)", 13, CYAN, True),
])
t(s, Inches(0.8), Inches(3.55), Inches(5.7), Inches(0.3),
  "0x0685C487Df4Cc0723Aa828C299686798294E9803", size=10, color=DIM, font="Consolas")

mt(s, Inches(0.6), Inches(3.95), Inches(6.0), Inches(0.5), [
    ("第 4 步  ·  Mint 1000 USDC 到 agent 钱包", 13, CYAN, True),
])

mt(s, Inches(0.6), Inches(4.4), Inches(6.0), Inches(0.5), [
    ("第 5 步  ·  EIP-712 签名 + Facilitator 上链", 13, CYAN, True),
])
t(s, Inches(0.8), Inches(4.75), Inches(5.7), Inches(0.3),
  "Python ref impl + TypeScript prod impl  ·  独立两套，相同链上效果", size=10, color=DIM)

# 链上 4 笔 tx 列表
card(s, Inches(0.6), Inches(5.2), Inches(6.0), Inches(1.5), BG_CARD, OAP_GRN)
t(s, Inches(0.8), Inches(5.3), Inches(5.7), Inches(0.3),
  "✅ 链上事实 (Blockscout immutable 可查)", size=11, color=OAP_GRN, bold=True)
mt(s, Inches(0.8), Inches(5.65), Inches(5.7), Inches(1.0), [
    ("•  合约部署 tx  0xb9bdfdb1...", 9, FG, False),
    ("•  Python e2e tx  0xff8a175e...", 9, FG, False),
    ("•  TypeScript e2e tx  0x5c10e2ae...", 9, FG, False),
    ("•  AWS Lambda tx  0x4562d26e...  (最新生产)", 9, OAP_GRN, True),
])

# 右半：HashKey faucet 截图
img_path = str(IMG_DIR / "hashkey-faucet.png")
try:
    s.shapes.add_picture(img_path, Inches(7.0), Inches(1.7), Inches(5.7), Inches(4.0))
    t(s, Inches(7.0), Inches(5.8), Inches(5.7), Inches(0.4),
      "↑ HashKey Testnet faucet  (Step 2 实操截图)", size=10, color=DIM, align="center")
except Exception as e:
    print(f"Image embed failed: {e}")

# 底部金句
card(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4), BG_CARD, GOLD)
t(s, Inches(0.8), Inches(6.9), Inches(12), Inches(0.3),
  "Python + TypeScript 两套独立实现产生完全相同的链上效果  →  证明协议层抽象正确",
  size=12, color=GOLD, bold=True, align="center")

n(s, """【HashKey 链上 demo】这是我做这个 demo 的完整 5 步——任何同事接 HashKey Chain 都可以照搬。

第 1 步：用 viem 生成 EVM 钱包
私钥本地生成，agent 地址 0x863d9C87b6bBd4aEf115C297C41643a0b887eAd7。

第 2 步：去 HashKey Testnet faucet 领测试币
URL: https://faucet.hashkeychain.net/faucet
右边这张截图就是 faucet 页面——粘贴你的钱包地址，得到约 1 HSK 测试币（用作 gas）。这是免费的，HashKey 提供。

第 3 步：用 Solidity 部署 MockUSDC 合约
这是 EIP-3009 的完整实现：transferWithAuthorization + cancelAuthorization + receiveWithAuthorization 三个核心函数。164 行 Solidity 代码。
合约地址：0x0685C487Df4Cc0723Aa828C299686798294E9803

第 4 步：mint 1000 USDC 到 agent 钱包

第 5 步：写 viem 客户端做 EIP-712 签名 + 上链
这一步我做了**两套独立实现**——Python 参考实现 + TypeScript 生产实现。

链上事实证据（左下卡片）：
- 合约部署 tx 0xb9bdfdb1...
- Python e2e tx 0xff8a175e... ✅
- TypeScript e2e tx 0x5c10e2ae... ✅
- **AWS Lambda 最新生产 tx 0x4562d26e...**（这是从今晚部署的 Lambda 上链的）

底部金句最关键：**Python + TypeScript 两套独立实现产生完全相同的链上效果——证明协议层抽象正确**。这是给 reviewer 的杀招——不是一套实现碰巧 work，是两套独立 implementation 都验证了同样的协议规范。

【DEMO 演示指引】如果有人问 "这是真的吗"——直接打开 testnet-explorer.hsk.xyz 输入合约地址 0x0685C487...，能看到所有这些 tx 的链上证据。链上数据 immutable，演讲后给老板审核也可以。""")
ftr(s, 15)

# ============================================================================
# Slide 16 — 三种集成形态 + CTA
# ============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "三种集成形态  ·  让所有客户都能用 AgentCore Payments",
    "OpenAgentPay 不卡死客户的成熟度路径  ·  Standalone → Plugin → AgentCore Native", OAP_GRN)

# 三个并排卡片
# Card 1: Standalone
card(s, Inches(0.6), Inches(1.8), Inches(4.0), Inches(3.8), BG_CARD, CYAN)
t(s, Inches(0.85), Inches(1.95), Inches(3.7), Inches(0.4), "形态 1", size=12, color=CYAN, bold=True)
t(s, Inches(0.85), Inches(2.35), Inches(3.7), Inches(0.6), "Standalone", size=22, color=FG, bold=True)
t(s, Inches(0.85), Inches(2.85), Inches(3.7), Inches(0.4),
  "✅ 当前 live", size=11, color=OAP_GRN, bold=True)
mt(s, Inches(0.85), Inches(3.3), Inches(3.7), Inches(2.2), [
    ("独立运行", 12, FG, True),
    ("CloudFront + Lambda", 11, DIM, False),
    ("不依赖 Bedrock", 11, DIM, False),
    ("", 6, DIM, False),
    ("适用客户", 11, CYAN, True),
    ("Web3 创业公司", 10, DIM, False),
    ("HashKey 自家 Agent", 10, DIM, False),
    ("不上 AWS Bedrock 的客户", 10, DIM, False),
])

# Card 2: Plugin
card(s, Inches(4.7), Inches(1.8), Inches(4.0), Inches(3.8), BG_CARD, GOLD)
t(s, Inches(4.95), Inches(1.95), Inches(3.7), Inches(0.4), "形态 2", size=12, color=GOLD, bold=True)
t(s, Inches(4.95), Inches(2.35), Inches(3.7), Inches(0.6), "Strands Plugin", size=22, color=FG, bold=True)
t(s, Inches(4.95), Inches(2.85), Inches(3.7), Inches(0.4),
  "🚧 下一步 1-2 周", size=11, color=GOLD, bold=True)
mt(s, Inches(4.95), Inches(3.3), Inches(3.7), Inches(2.2), [
    ("轻量集成", 12, FG, True),
    ("AgentCore Runtime", 11, DIM, False),
    ("AgentCore Identity", 11, DIM, False),
    ("AgentCore Memory", 11, DIM, False),
    ("Observability", 11, DIM, False),
    ("", 6, DIM, False),
    ("适用客户", 11, GOLD, True),
    ("Bedrock + Strands 客户", 10, DIM, False),
])

# Card 3: AgentCore Native
card(s, Inches(8.8), Inches(1.8), Inches(4.0), Inches(3.8), BG_CARD, OAP_GRN)
t(s, Inches(9.05), Inches(1.95), Inches(3.7), Inches(0.4), "形态 3", size=12, color=OAP_GRN, bold=True)
t(s, Inches(9.05), Inches(2.35), Inches(3.7), Inches(0.6), "AgentCore Native", size=22, color=FG, bold=True)
t(s, Inches(9.05), Inches(2.85), Inches(3.7), Inches(0.4),
  "⏳ 等 AWS 开放 BYO", size=11, color=PURPLE, bold=True)
mt(s, Inches(9.05), Inches(3.3), Inches(3.7), Inches(2.2), [
    ("深度集成", 12, FG, True),
    ("BYO connector", 11, DIM, False),
    ("AWS Console 内置", 11, DIM, False),
    ("业务代码 0 改动", 11, DIM, False),
    ("", 6, DIM, False),
    ("适用客户", 11, OAP_GRN, True),
    ("全部 AWS 客户", 10, DIM, False),
    ("(主流市场)", 10, DIM, False),
])

# 底部 CTA
card(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.0), BG_CARD, OAP_GRN)
t(s, Inches(0.8), Inches(5.95), Inches(12), Inches(0.4),
  "我们一起做  ·  让所有 AWS 亚太客户都能用 Payments",
  size=18, color=OAP_GRN, bold=True, align="center")
mt(s, Inches(0.8), Inches(6.35), Inches(12), Inches(0.5), [
    ("github.com/neosun100/openAgentPay  ·  Apache 2.0  ·  欢迎 PR / Issue / 一起接 connector", 12, FG, False),
])

n(s, """【三种集成形态 + CTA】最后总结。

OpenAgentPay 设计了三种集成形态，让客户根据自己情况选——这是商业化的杀招，**任何客户都能 self-identify 自己是哪种**。

**形态 1：Standalone（当前 live）**
独立 Lambda + CloudFront + API Gateway 部署。客户即使不用 Bedrock 也能用。
适用：Web3 创业公司、HashKey 自家 Agent 平台、不上 AWS Bedrock 的客户。

**形态 2：Strands Plugin（下一步 1-2 周）**
作为 Strands Plugin 接入 AgentCore Runtime + Identity + Memory + Observability。业务代码 1 行不改：把 AgentCorePaymentsPlugin 换成 OpenAgentPayPlugin。
适用：用 Bedrock + Strands 的客户（最大客户群）。

**形态 3：AgentCore Native（等 AWS 开放 BYO connector）**
当 AWS 开放 BYO connector 接口，我们直接 register 进去。客户在 AgentCore Console 里可以下拉选 HashKey Chain、Binance Pay、OKX Pay。业务代码 0 改动。
适用：全部 AWS 客户（主流市场）。

**Call to Action（最后 3 件事）**：

第一，我做这个不是单兵作战。目标是让我们 AWS Web3 团队作为 AWS 在亚洲的 **Agent Payments 解决方案中心**。我邀请大家一起加入：

第二，**谁的客户有 Payments 需求？欢迎找我聊**。
- HashKey、OKX、Bitget、Bybit 任何亚洲交易所
- 支付宝、微信、UnionPay 任何传统支付
- MetaMask、WalletConnect 任何 Web3 自托管

按 connector 模板接，1-2 天就能给客户 ship 一个 PoC。

第三，**有兴趣一起做 SA 协作的同事**，可以一起负责一个钱包/协议方向。代码全开源 Apache 2.0，github.com/neosun100/openAgentPay。

终极意义：让所有 AWS 客户都能用 Agent Payments，不只是北美 Coinbase + Stripe 用户。这是 AWS Web3 SA 团队该做的事。

谢谢大家。Q&A 时间——任何问题、想法、批评都欢迎。""")
ftr(s, 16)

# ============================================================================
#  SAVE
# ============================================================================
out = Path(__file__).parent / "openagentpay-talk.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
print(f"   Total slides: {len(prs.slides)}")

# Verify all notes
total_chars = 0
for i, slide in enumerate(prs.slides, 1):
    notes = slide.notes_slide.notes_text_frame.text
    chars = len(notes)
    total_chars += chars
    title_text = "(no title)"
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            title_text = shape.text_frame.text.strip().split("\n")[0][:50]
            break
    print(f"   Slide {i:2d}  notes={chars:4d}  {title_text}")
print(f"\n   📊 Total notes: {total_chars} chars  ({total_chars / len(prs.slides):.0f} avg/slide)")
