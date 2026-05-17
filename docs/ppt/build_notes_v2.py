"""
build_notes_v2.py — 重写所有演讲备注

读 openagentpay-talk.pptx，全部 16 页备注重写：
  - 立场调整：从"宣讲我做的牛 X 项目"→"抛砖引玉，邀请兄弟们一起完善"
  - 称呼：用"兄弟们"（Web3 团队内部分享自然语气）
  - 透彻讲清楚 8 个 challenge 问题
  - HSK vs USDC 真相
  - OAP-CEX 命名 + 含义
  - 5 层架构类比
  - CEX 不上链原理
  - Demo 演讲台操作流程

不动 slide 视觉，只改 notes_text_frame。
"""
from pathlib import Path
from pptx import Presentation

PPT_PATH = Path(__file__).parent / "openagentpay-talk.pptx"
prs = Presentation(str(PPT_PATH))

print(f"读入: {PPT_PATH.name}  ·  {len(prs.slides)} 页")

# ============================================================================
#  备注内容（按 slide 编号）
# ============================================================================
NOTES = {}

# ----------------------------------------------------------------------------
# Slide 1 — 封面
# ----------------------------------------------------------------------------
NOTES[1] = """【封面 · 开场白】

兄弟们好。今天这个分享是受周会安排来的，主题是 AgentCore Payments + OpenAgentPay。

简单说一下背景。上周 SSO 给我们团队讲了 AgentCore Payments 这个上周 AWS 刚发布的新产品。我个人对它做了一些深入研究，写了一篇公众号文章解读它，但是越研究越发现一个问题——这个产品当前 Preview 阶段只支持 Coinbase 和 Stripe 两个钱包，对我们 Web3 团队覆盖的亚太合规交易所、CEX 钱包、Web3 自托管钱包的客户，**全部用不上**。

所以我自己做了一些扩展，叫 OpenAgentPay。当前还是 work-in-progress 状态，但有了一些初步成果。今天来分享给兄弟们，**主要目的不是说这个东西多好——是想抛砖引玉，邀请大家一起完善**。

为什么这么说？因为这个事说实话**单兵作战做不成**。我们覆盖那么多种类型的客户、那么多个交易所、那么多种支付方式——一个人一辈子也接不完。我做了个起点，邀请兄弟们一起把后面的接进来，把这个事做成我们 FSI DNB SA Team 整体的贡献。

整个分享大概 15-20 分钟。前半我讲 AgentCore Payments 是什么、它的短板；后半讲我做的 OpenAgentPay 怎么补这个短板。任何时候**有想法、有挑战、有反对意见都欢迎打断**——这正是我希望听到的。"""

# ----------------------------------------------------------------------------
# Slide 2 — 目录
# ----------------------------------------------------------------------------
NOTES[2] = """【目录】

今天的分享分两部分。

**Part 1：AgentCore Payments**
3 个内容：场景钩子（Agent 为什么需要钱包）、它怎么工作（x402 协议）、它的现状和短板（这是关键 — 短板是后半的引子）。

**Part 2：OpenAgentPay**
3 个内容：解决方案设计（5 层架构）、HashKey 实战（demo 看链上效果）、怎么一起做（邀请兄弟们参与）。

我特别想强调一点——**第一部分讲 AgentCore Payments 的"短板"不是为了 trash 这个产品**。恰恰相反，AgentCore Payments 是个非常 promising 的好产品，AWS 在 Agent 经济卡位的关键一步。但任何 Preview 产品都有覆盖面有限的问题，AWS 也不可能一上来就支持所有场景。

我们 Web3 团队的角度看，**短板就是机会**——AWS 没接的钱包、没覆盖的协议，正是我们能贡献的地方。

兄弟们听的时候记一个心态——**这是个 community effort，不是某个人的项目**。我做了起点，但后面长成什么样、跑得多远，得看在座兄弟们一起怎么推。"""

# ----------------------------------------------------------------------------
# Slide 3 — 章节分隔（第 1 章）
# ----------------------------------------------------------------------------
NOTES[3] = """【第一章 · 什么是 AgentCore Payments】

我们先看这个产品本身。

AgentCore Payments 是 AWS 上周（2026 年 5 月 7 日）金融服务峰会上发布的，会被未来定义为"Agent 经济元年的里程碑"——因为它是首个**面向自主 AI Agent 的托管支付基础设施**。

接下来 4 页我从场景钩子开始，把它讲清楚。这部分理论性强一点，但是后半部分的基础——理解了这部分，才能看懂我们做的 OpenAgentPay 在解决什么问题。"""

# ----------------------------------------------------------------------------
# Slide 4 — 场景钩子
# ----------------------------------------------------------------------------
NOTES[4] = """【场景钩子 · Agent 能自己付款吗】

兄弟们想象一个场景。你的 AI Agent 在帮你做金融研究，分析到一半发现需要访问彭博终端的付费数据才能给出更准确的结论。这个时候 Agent 能自己付款吗？

**在今天之前，答案是不能**。Agent 必须停下来等人类介入：开账号、绑信用卡、管理 API Key、处理订阅...走完一整套人类世界为人类设计的支付流程。

这就是 AI Agent 领域一个被所有人忽视但致命的矛盾——**智能是自主的，付款不是**。

为什么这是个大问题？因为 Agent 经济正在快速到来。Anthropic 的 Computer Use、OpenAI 的 Operator、AWS 自己的 Strands Agent，都在让 Agent 能自主使用工具。但是只要付款这一环还要人参与，Agent 的自主性就被打了折扣——只能做"调用免费 API"，不能做"研究决策时按需付费查数据"这种真正有商业价值的任务。

AWS 的 AgentCore Payments 就是要解决这个问题——**给 AI Agent 一张能自己用的钱包**。这是 Agent 经济的关键基础设施层。

下一页讲它具体是什么。"""

# ----------------------------------------------------------------------------
# Slide 5 — AgentCore 一句话定义
# ----------------------------------------------------------------------------
NOTES[5] = """【AgentCore Payments 是什么】

一句话定义：**首个面向自主 AI Agent 的托管支付基础设施**——让 Agent 在推理循环内自主发现、认证、结算、获取付费资源，无需人类介入。

注意三个关键词：

**"首个"**——微软 Azure 和 Google Cloud 在 Agent 平台层都没有同类产品。AWS 在 Agent 经济领先一代。

**"托管"**——AWS 一手包办协议协商、钱包管理、限额控制、合规检查、可观测性。开发者只需要几行代码就能接入，不用自己写复杂的钱包逻辑。

**"推理循环内"**——这是最关键的设计——付款动作不打断 Agent 思考，就像人喘气一样自然。Agent 调用工具的时候碰到 402 Payment Required，AgentCore 自动处理付款，Agent 完全无感。

看这张图，AgentCore Payments 不是独立产品，而是 AgentCore 平台的**原生模块**，跟 Identity、Gateway、Observability、Memory、Runtime 平起平坐。这就是 AWS 的设计哲学——'native, not bolted-on'，原生集成不是贴片。意味着企业安全团队审批一次就行，不用重新过评审。

它包含 4 个子能力：
- **Payment Manager**：编排大脑（协议协商 + 交易签名 + 结算）
- **Payment Guardrail**：风控边界（授权 + 限额 + 预算强制）
- **Wallet Providers**：钱包对接层（Coinbase CDP / Stripe Privy）
- **Protocols**：协议层（当前 x402，未来可扩展）

下一页讲它的核心协议 x402 怎么工作。"""

# ----------------------------------------------------------------------------
# Slide 6 — x402 协议
# ----------------------------------------------------------------------------
NOTES[6] = """【x402 协议 · HTTP 402 状态码】

要理解 AgentCore Payments，必须先懂 x402。x402 是 Coinbase 在 2024 年底发布的开放协议，做了一件精妙的事——**把 HTTP 协议里那个一直没人用的状态码 402 Payment Required 激活了**。

5 步搞定一次 Agent 付款：

第 1 步：Agent 发起请求 GET /resource，普通 HTTP 调用。
第 2 步：服务端返回 402 Payment Required，告诉你要付多少、付给谁、用啥币。
第 3 步：Agent 内部，AgentCore Payments 查预算 → 拿钱包私钥 → 做 EIP-3009 离线签名。
第 4 步：Agent 带上 X-Payment 头重试请求。
第 5 步：服务端验证签名 → Facilitator 上链结算 → 返回 200 OK + 你要的内容。

整个流程，**Agent 不需要持有 gas、不需要连区块链节点、不管理 nonce**。复杂度全部被 x402 Facilitator 吸收掉了，Agent 的 UX 跟调 REST API 一模一样。

底部这一组数字告诉兄弟们 x402 不是玩具：30 天 7541 万笔交易、累计 1.69 亿笔、活跃买家 59 万、卖家 10 万、L2 结算约 2 秒完成 finalize。

被 Stripe、AWS、Cloudflare、Vercel、Quicknode、Messari、Alchemy 都信赖。**AWS 选 x402 不是保守的选择——是押最成熟、已生产验证、增长最快的赛道**。这是个深思熟虑的战略选型。

【备注】其实当前 Agent Payments 协议格局有 10+ 个：x402（Coinbase）、MPP（Stripe + Tempo IETF Draft）、AP2（a16z 推）、ACP（多家 startup）、TAP、UCP 等。AWS 选 x402 是因为它在生产数据、生态成熟度、技术成熟度三个维度都最强。"""

# ----------------------------------------------------------------------------
# Slide 7 — 经济学
# ----------------------------------------------------------------------------
NOTES[7] = """【经济学 · 3000x 差距】

兄弟们可能问：让 Agent 自动付款，直接绑张信用卡不就完了？答案是：经济学不允许。

看这张图：**传统 Stripe 信用卡单笔最低手续费 0.30 美元，x402 单笔成本约 0.0001 美元——差了整整 3000 倍**。

举个具体例子：Agent 调用一次 API 可能只值 0.001 美元，也就是 0.1 分钱。但走传统信用卡要扣 0.30 美元，30 分钱手续费。**付款比买的东西贵了 300 倍**——这个市场根本上不了规模。

所有按次计费的 Pay-per-use 市场——API 调用、AI 推理调用、按字数计费的内容、按次的数据查询——都被传统支付的单位经济学**彻底锁死**几十年。Agent 经济需要一个新支付层。

而 x402 + 稳定币 USDC on Base L2 解锁了这个市场：

- **L2 结算 ~2 秒**（Flashblock 预确认 200 ms），传统 ACH/Card 要 T+2-3 天
- **单笔成本 < 1 分钱**，让小于 1 美元的微支付变得理所当然
- **零账户设置**——有钱包就能付，不用注册商户
- **无 API Key 风险**——签名是一次性授权，泄露也无害
- **天然全球同价**——不管在哪个国家，费率都一样

让小于 1 美元的 Agent 微支付**从不可能变成理所当然**。这就是 AWS、Stripe、Cloudflare 同时押注这个协议的真正原因——**不是因为它酷，而是因为它解锁了一个被锁住几十年的万亿级市场**。

好，到这里 AgentCore Payments 该讲的部分讲完了。下一章讲它的现状和短板——这是我们 OpenAgentPay 工作的起点。"""

# ----------------------------------------------------------------------------
# Slide 8 — 章节 2 分隔
# ----------------------------------------------------------------------------
NOTES[8] = """【第 2 章 · 现状与短板】

刚才讲的 AgentCore Payments 听起来全是优点。但任何 Preview 产品都有覆盖面有限的问题，AWS 也不可能一上来就支持所有场景。

接下来 2 页，我把数据摆出来——**不是为了 trash 这个产品，而是为了看清楚我们的位置：哪些客户能用、哪些不能用、机会在哪里**。

兄弟们听的时候特别注意——**短板不是产品的过错，是我们 Web3 团队的机会**。"""

# ----------------------------------------------------------------------------
# Slide 9 — 现状 + 短板
# ----------------------------------------------------------------------------
NOTES[9] = """【现状与短板】

这页把数据摆出来。

**左边是 x402 生态的现状——非常成熟**。30 天 7541 万笔交易、活跃买家 59 万、卖家 10 万、单笔成本不到 1 分钱。这是一个生产级、被头部玩家信赖的协议层——AWS 押得对。

**右边这张图是 AgentCore Payments 当前支持的两种钱包路径**：

第一种 Coinbase CDP——Custodial Wallet 模式，钱包私钥由 Coinbase 托管。优势是 fast onboarding，缺点是用户依赖 Coinbase。
第二种 Stripe Privy——Embedded Wallet 模式，每个用户一个 EOA 钱包，私钥分片在用户设备 + Privy 后端。

看似覆盖很广——一个 custodial、一个 self-custody。**但仔细看：**

**Coinbase 是美国上市公司**，主要服务北美 Web3 用户。
**Stripe Privy 是 2024 年 Stripe 收购的**，主要服务北美和欧洲合规支付场景。

下面这 4 条限制特别关键：

1. **仅 2 个钱包**
2. **仅 1 个协议（x402）**——MPP / AP2 / ACP 等正在涌现的 Agent 支付协议都没接
3. **仅 4 个 region 可用**（us-east-1, us-west-2, eu-west-1, ap-southeast-2）——**注意，没有亚洲 region**。新加坡、东京、香港、孟买，AgentCore Payments 都用不了
4. **Preview 阶段，BYO connector 接口未开放**——意味着第三方钱包没法 register 进去

兄弟们到这里应该有感觉——这个产品在亚洲场景、CEX 场景，**当前都不能用**。下一页讲我们的客户具体是哪些。"""

# ----------------------------------------------------------------------------
# Slide 10 — 转折页
# ----------------------------------------------------------------------------
NOTES[10] = """【那个关键问题 · 我们的客户在哪里】

所以问题来了：我们 AWS Web3 团队，服务的客户群体——亚洲合规交易所、亚洲 CEX Pay、Web3 自托管钱包、传统亚洲支付——**4 个类别全部当前用不上 AgentCore Payments**。

这不是 AWS 产品做得不好。Coinbase + Stripe 已经覆盖了北美和欧洲的主要场景。这是 AWS 推产品的策略选择——先做最大公分母，等市场反馈再扩张。

但这就把我们 Web3 团队卡死了。

我们去找 HashKey 推 AgentCore Payments，HashKey 会问："你们支持 HashKey Chain 吗？支持 HKDR 港币稳定币吗？"答：暂时不支持。
找 Binance 推，会问："支持 Binance Pay API 吗？"答：不支持。
OKX、Bitget、Bybit、支付宝、微信......同样的问题。

我们 Web3 团队的核心 KPI 是覆盖亚太 Web3 客户。**手里这个产品，目标客户一个都用不上**。

兄弟们，我们怎么办？

**选项 1：等 AWS 自己开放支持**。但 roadmap 写的是 'Others* coming soon'——什么时候、是否包含 HashKey、AWS 没给时间表，我们等不起。

**选项 2：放弃这个产品**，跟客户说 AWS 这个东西好但你用不上。但作为 AWS SA，这个 narrative 太弱了——客户会觉得我们不专业。

**选项 3：自己做扩展层**，让我们的客户也能用。**这是我选的路径——OpenAgentPay**。

但**单兵作战做不成**——所以今天来邀请兄弟们一起做。接下来 6 页讲清楚 OpenAgentPay 是什么、当前到哪里、未来兄弟们能怎么参与。"""

# ----------------------------------------------------------------------------
# Slide 11 — 章节 3 分隔
# ----------------------------------------------------------------------------
NOTES[11] = """【第 3 章 · OpenAgentPay 解决方案】

接下来 5 页讲 OpenAgentPay。

**请兄弟们注意一个心态**——这是 work-in-progress，不是 production，不是 done deal。我会诚实地告诉哪里完成了、哪里还在开发、哪里需要兄弟们加入。

但当前的进度已经有一些**有真东西的成果**：协议层验证完成、链上 e2e 跑通、AWS 部署 live、4 笔真上链 tx 在 Blockscout 永久可查。

**我说"邀请兄弟们参与"不是空话**——后面有具体的参与方式。"""

# ----------------------------------------------------------------------------
# Slide 12 — 5 层架构（关键讲解）
# ----------------------------------------------------------------------------
NOTES[12] = """【5 层可插拔架构】

OpenAgentPay 一句话定位：**AgentCore Payments 的开放可插拔扩展层**——让任何钱包、任何协议、任何稳定币都能即插即用接入。

我用一个**餐厅类比**让兄弟们一秒理解 5 层架构：

- Layer 1 · Strands Plugin → **服务员的点单 PAD**（Agent 直接用的 SDK）
- Layer 2 · Payment Orchestrator → **餐厅经理**（编排订单、检查预算、协议路由）
- Layer 3 · Protocols → **支付方式**（信用卡、微信、现金，可选）
- Layer 4 · Wallet Connectors → **POS 机**（每种支付方式一台，按统一接口）
- Layer 5 · Self-Hosted Facilitator → **银行后台清算系统**

**这 5 层的核心设计原则**：Layer 2/3/5 是 framework 不变；Layer 1 是 SDK 入口不变；**只有 Layer 4 按钱包变化**——每个钱包写一个 connector，按统一接口。

**用具体例子说明这个设计的价值**：

假如今天 HashKey 工程师想给他的 Agent 加付款能力：

❌ 不用 OpenAgentPay：自己实现协议 + 钱包 + 上链 + 错误处理 + 安全 — **估时 1-2 个月 + 1 个工程师**

✅ 用 OpenAgentPay：在 Layer 4 写一个 HashKeyChainConnector（参考已有模板）— **估时 1-2 天 + 1 个工程师**

新加入 OKX？同样在 Layer 4 写一个 OKXConnector，**1-2 天**。其他层不动。

类比 Kubernetes：CRI（Container Runtime Interface）/ CSI / CNI 让任何容器、存储、网络可插拔接入。**OpenAgentPay 想做 Agent Payments 的 CRI 时刻**——定义 WalletConnector + ProtocolAdapter 标准接口。

底部那行金色字是核心承诺：**业务代码改 1 行 → 切换钱包**。
- payment_manager = PaymentManager(wallet_provider="hashkey-chain")
- payment_manager = PaymentManager(wallet_provider="binance-pay")
- payment_manager = PaymentManager(wallet_provider="coinbase-cdp")

这就是 framework 的真正价值——**让客户不锁死技术栈**。

兄弟们如果想参与的话——主要是 Layer 4。每个人选一个钱包做 connector，1-2 天 ship 一个 demo，给客户立刻能用。这就是我说"community effort"的具体形态。"""

# ----------------------------------------------------------------------------
# Slide 13 — 双协议轨道（关键讲解）
# ----------------------------------------------------------------------------
NOTES[13] = """【双协议轨道 · OAP-CEX 是什么】

这页讲一个深刻的问题——这是项目最容易被 challenge 的核心 architectural question：**为什么 Binance 不直接用 x402？**

**先说 OAP-CEX 名字怎么来的**：

OAP = OpenAgentPay 缩写
CEX = Centralized Exchange（中心化交易所）

合起来 = **OpenAgentPay 项目下专门给 CEX 用的协议**。命名借鉴 x.509 / x402 / TLS 等工业界传统——短前缀 + 用途后缀。未来扩展可以叫 OAP-FIAT（传统支付）、OAP-BANK（银行直连）。

**回到核心问题：为什么 Binance 不能用 x402？**

简短回答：**x402 是链上协议，但 Binance 这类 CEX 结构上不上链**。

兄弟们想想——你在 Binance 看到的"1000 USDT 余额"，**不是链上 USDT 合约的状态，是 Binance 内部数据库里的一行记录**：user_id=12345, asset=USDT, balance=1000。

CEX 的真相：**99% 的资金流动在内部账本，1% 在链上（提现时才上链）**。这是 CEX 的核心架构，跟"区块链"思路根本不一样：

- 内部转账：Binance 数据库 update，~50ms 完成，**0 gas**
- 提现到链：才上链，要 gas，几秒确认

**为什么 CEX 这么设计？**

- 成本：链上转账要 gas，CEX 内部记账 0 成本
- 速度：链确认要几秒，内部记账毫秒级
- 合规：KYC 信息只能在 CEX 内部，链上是 pseudonymous
- 效率：高频交易用户每天几百笔，全上链链就堵了

**所以强行套 x402 给 Binance 等于让中心化数据库假装成 ERC-20 合约——既笨拙又违背 CEX 的低成本优势。**

**OpenAgentPay 的关键 insight**：x402 协议的"形状"很好（402 challenge → sign → retry），但"加密层 + 结算层"应该可插拔。所以我们：

- **x402** = 协议形状 + EIP-712 签名层（链上钱包用）
- **OAP-CEX** = 同样协议形状 + HMAC 签名层（CEX 用）

两者**共享 ProtocolAdapter 接口**；PaymentManager 通过 ProtocolRouter 自动按 402 response signature 派发。

**为什么客户会选 OAP-CEX？**

- HashKey、Binance、OKX 客户：他们的用户**已经在 CEX**，不需要折腾链
- 传统支付场景（支付宝、微信、Stripe credit card）：根本不上链
- B2B 高额场景：KYC 必需，CEX 路径更合规

**给兄弟们一个杀招总结**：

> "x402 让 Web3 user 能付款。OAP-CEX 让剩下 90% 的用户也能付款。世界上每天 99% 的支付不在链上——OAP-CEX 是把 Agent Payments 从 Web3 niche 推向 mainstream 的桥梁。"

这是 OpenAgentPay 跟 AgentCore 的真正差异化——AgentCore 只走链上路径，**我们双轨**。"""

# ----------------------------------------------------------------------------
# Slide 14 — Live AWS（详细链路 + demo 演示指引）
# ----------------------------------------------------------------------------
NOTES[14] = """【Live on AWS · 完整合规架构】

OpenAgentPay 不是 ppt-only project——demo 已经部署到 AWS 生产环境。

**完整链路（按数据流）**：

1. 用户浏览器 → HTTPS POST /api/pay
2. CloudFront（CDN + DDoS + TLS + CORS）→ /api/* 路由到 API Gateway
3. API Gateway HTTP API（合规公网入口）→ 转 Lambda integration
4. Lambda（Node 20）→ 查 Secrets Manager 拿私钥 → viem 构造 EIP-712 签名
5. Lambda → HashKey Chain RPC broadcast → 等 ~5 秒确认
6. 返回 tx hash → API Gateway → CloudFront → 用户浏览器显示 ✅

**为什么走这个链路（不能简化）**：

- CloudFront 不能去掉：Lambda 直接公开会被攻击 + 没 CDN
- API Gateway 不能去掉：用 Lambda Function URL 会被 Palisade 检测 mitigation（这是真的——昨晚我们就遇到了）
- Lambda 不能去掉：私钥不能放浏览器
- Secrets Manager 不能去掉：私钥必须 KMS 加密静态存储

**4 个合规原则**：
1. No Function URL（Palisade-proof）
2. IAM-scoped throughout
3. KMS at rest
4. Secrets never in logs

这套架构跟 AWS AgentCore Payments 自己用的是**同一套模式**——客户接进来零学习成本。

【Performance】
- Cold start ~1s
- End-to-end payment ~5s
- CDK deploy 357s
- 4 verified on-chain tx
- 82 tests pass

══════════════════════════════════════════
【🎬 DEMO 演讲台操作流程】（建议 3-4 分钟）
══════════════════════════════════════════

**这一页停留时切到浏览器：https://d1p7yxa99nxaye.cloudfront.net**

Step 1：在 Tab 1 'Run Demo' 停留——给兄弟们看 4 步流程概览
说："这是个三 Tab 的 demo，先看 Tab 1 怎么手动跑链上结算。"

Step 2：点 Step 1 'Run' button
说："这调 /api/wallet 查链上 USDC 余额。等 1 秒——你看，~998 USDC，是从 HashKey Chain 真链上读的。"
（如果有同事问"这 USDC 哪来的"——这就是 next slide 要讲的，先 hold）

Step 3：点 Step 2 'Run'
说："创建 Payment Session，预算 $1，TTL 60 分钟。**这是 spend governor 边界**——超预算硬拒绝，不依赖 LLM 判断，基础设施层强制。"

Step 4：点 Step 3 'Pay'
说："这是核心步骤——EIP-712 typed data 签名 + Facilitator 上链。等 5 秒。"
（5 秒后看到 tx hash 返回）
说："看，tx 0x...拿到了。点这个 Blockscout 链接——"
**点击链接**，浏览器跳转
说："这就是真上链的证据。block 号、gas、payer、recipient——immutable 永远可查。任何同事会后想验证，自己查 testnet-explorer.hsk.xyz 输合约地址 0x0685C487 就能看到。"

Step 5（如果时间够）：切 Tab 3 'AI Agent'
说："这个 Tab 模拟 Strands Agent 自主决策——免费工具直接用，付费工具触发链上结算。下一步会接通真 Bedrock + Strands。"
点"付费 ETH 深度分析"按钮 → 等 5 秒 → 看到 tx hash
说："看，又是一笔真上链。这就是 Agent 自主付费的最终形态。"

**Demo 应急方案**：
- 如果 demo 卡了 → 切到 https://github.com/neosun100/openAgentPay 截图展示
- 如果 RPC 卡 → "HashKey testnet 节点偶尔会卡，是 testnet 问题，协议层是稳定的"
- 完全 down → 用 Blockscout 历史 tx 链接展示链上事实

══════════════════════════════════════════"""

# ----------------------------------------------------------------------------
# Slide 15 — HashKey 链上 demo（HSK + USDC 故事）
# ----------------------------------------------------------------------------
NOTES[15] = """【HashKey Chain 链上 demo · HSK 和 USDC 真相】

这是我做这个 demo 的完整 5 步——任何兄弟接 HashKey Chain 都可以照搬。

**特别注意一个细节**——很多兄弟会问：你从 faucet 领的是 1 个 HSK，但 demo 里余额显示 1000 USDC，这是怎么回事？

══════════════════════════════════════════
【HSK 和 USDC 的真相】（你自己也会被问）
══════════════════════════════════════════

**HSK 和 USDC 是两个完全不同的代币，扮演两个完全不同的角色，跟以太坊上 ETH 和 USDT 关系一模一样。**

类比以太坊：
- ETH（原生币，付 gas）  ←→  USDT/USDC（ERC-20，业务转账）

类比 HashKey Chain：
- HSK（原生币，付 gas）  ←→  MockUSDC（我部署的 ERC-20，业务结算）

**我做了什么（演讲台讲清楚）**：

第 1 步：我去 HashKey Testnet faucet 领币。
faucet URL：https://faucet.hashkeychain.net/faucet
领到的就是这 1 个 HSK，截图里看到的就是它。
**用途：作为"汽油"，每次链上交易要烧一点 HSK 当 gas**。

第 2 步：但 1 个 HSK 不能做演示金额——HSK 价格不稳定，业务金额需要稳定币（1 USDC = 1 USD）。

第 3 步：HashKey Chain Testnet 上**没有官方 USDC**——Circle 还没在 HashKey Chain 部署，这是 HashKey 路线图上的事。

第 4 步：所以我自己写了一个 MockUSDC 合约（ERC-20 + EIP-3009 完整实现）。用我那 1 HSK 当 gas，部署到 HashKey Chain Testnet。
合约地址：0x0685C487Df4Cc0723Aa828C299686798294E9803

第 5 步：调用 mint() 函数，给自己钱包打 1000 个 MockUSDC。因为 mock token 嘛，mint 函数我自己写的，没限制。

**现在钱包：约 1 HSK（剩一点 gas）+ 1000 USDC（自己 mint 的）**

每次 demo 付款：
- 业务转账：扣 0.001 USDC（演示给兄弟们看的"业务金额"）
- 链交易 gas：~0.0000001 HSK（极小，用户感知不到）

**关键 insight**：
USDC 是真上链的、签名是真的、合约逻辑是真的——只有 mint 这 1000 个是我自己 mint 的。等 Circle 在 HashKey Chain 部官方 USDC，**把合约地址换一下就是生产级，协议层 0 改动**。

这就是为什么这个 demo 是 production-equivalent 但不是 production-ready——把 mock 换成 Circle 官方就 ready。

══════════════════════════════════════════
【链上事实证据】（兄弟们会后自己查）
══════════════════════════════════════════

- 合约部署 tx: 0xb9bdfdb1...
- Mint tx: 0xd862e80e...
- Python e2e tx: 0xff8a175e... ✅
- TypeScript e2e tx: 0x5c10e2ae... ✅
- AWS Lambda 最新生产 tx: 0x4562d26e... ✅

**Python + TypeScript 两套独立实现产生完全相同的链上效果——证明协议层抽象正确**。

这是给 reviewer 的杀招——不是一套实现碰巧 work，是两套独立 implementation 都验证了同样的协议规范。任何兄弟会后想验证，testnet-explorer.hsk.xyz 输合约地址 0x0685C487 就能看到所有 tx。"""

# ----------------------------------------------------------------------------
# Slide 16 — CTA（邀请兄弟们参与的具体方式）
# ----------------------------------------------------------------------------
NOTES[16] = """【三种集成形态 + 邀请兄弟们一起做】

**先讲三种集成形态**——这是商业化的杀招，让任何客户都能 self-identify 自己是哪种。

**形态 1：Standalone（当前 live）**
独立 Lambda + CloudFront + API Gateway 部署。客户即使不用 Bedrock 也能用。
适用：Web3 创业公司、HashKey 自家 Agent 平台、不上 AWS Bedrock 的客户。

**形态 2：Strands Plugin（下一步 1-2 周）**
作为 Strands Plugin 接入 AgentCore Runtime + Identity + Memory + Observability。业务代码 1 行不改：把 AgentCorePaymentsPlugin 换成 OpenAgentPayPlugin。
适用：用 Bedrock + Strands 的客户（最大客户群）。

**形态 3：AgentCore Native（等 AWS 开放 BYO connector）**
当 AWS 开放 BYO connector 接口，我们直接 register 进去。客户在 AgentCore Console 里下拉选 HashKey Chain、Binance Pay、OKX Pay。业务代码 0 改动。
适用：全部 AWS 客户（主流市场）。

══════════════════════════════════════════
【真正的 CTA · 兄弟们一起做】
══════════════════════════════════════════

我做这个不是单兵作战。**兄弟们说实话，一个人接不完所有钱包**。我列一下当前需要兄弟们参与的具体方式：

**第一种参与方式：负责一个钱包/协议方向**

如果你覆盖某个交易所或钱包客户：
- 覆盖 OKX 客户的兄弟 → 一起做 OKXConnector，1-2 天
- 覆盖 Bitget 客户的兄弟 → 一起做 BitgetConnector
- 覆盖 Bybit 客户的兄弟 → 一起做 BybitConnector
- 覆盖 HashKey Pro（CEX 不是 Chain）→ HashKey Pro Sandbox connector
- Binance 已经我做了 work-in-progress

**第二种参与方式：客户带场景进来**

如果你的客户问起 Payments，找我聊。可能：
- 客户用 MetaMask → 1 天写个 MetaMaskConnector
- 客户用 Stripe → 1 天 + 1 天加传统支付协议轨道
- 客户在做 RWA → 把 OpenAgentPay 当成支付层 demo 给客户

**第三种参与方式：架构 review + 反馈**

OpenAgentPay 的协议层、framework、安全设计——兄弟们看了有想法直接提 issue / PR。开源 Apache 2.0，github.com/neosun100/openAgentPay。

**第四种参与方式：一起对外发声**

如果觉得 OK，我们可以把这个工作输出成 FSI DNB SA Team 的 community contribution——内部分享、外部博客、再到 AWS Solutions Library。这就是我们团队的 visibility。

══════════════════════════════════════════
【最后一句话】
══════════════════════════════════════════

兄弟们，**个人力量很渺小**——一个人接一个 connector 可能要 1-2 天，但接 20 个钱包就要 1-2 个月。**团队拧成一股绳，每人 1-2 天，2 周就能覆盖整个亚太市场**。

这就是为什么我说"邀请兄弟们一起做"——不是客套，是实打实的需求。

终极意义：**让所有 AWS 客户都能用 Agent Payments，不只是北美 Coinbase + Stripe 用户**。这是我们 AWS Web3 SA 团队该做的事，作为 FSI DNB SA Team 的整体贡献。

谢谢兄弟们。Q&A 时间——任何问题、想法、批评、challenge 都欢迎。

GitHub: github.com/neosun100/openAgentPay
Live demo: d1p7yxa99nxaye.cloudfront.net
私聊: jiasunm@amazon.com"""

# ============================================================================
#  应用所有 notes 到 PPT
# ============================================================================
for i, slide in enumerate(prs.slides, 1):
    if i in NOTES:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = NOTES[i]
        print(f"  ✓ Slide {i:2d}  notes={len(NOTES[i]):4d} chars")

# 保存
out_path = Path(__file__).parent / "openagentpay-talk.pptx"
prs.save(out_path)
print(f"\n✅ 已保存: {out_path}")

# 验证总字数
total_chars = 0
for i, slide in enumerate(prs.slides, 1):
    notes = slide.notes_slide.notes_text_frame.text
    total_chars += len(notes)
print(f"   📊 总演讲备注: {total_chars} 字  ({total_chars / len(prs.slides):.0f} 平均/页)")
