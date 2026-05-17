"""
build_full_v2.py — 完整重建 23 页 PPT（基于 deep-research-ppt skill）

变化:
  - 标题：AgentCore Payments 主标，OpenAgentPay 副标
  - 前半（章节 1）从 5 页 → 11 页（加 6 页：万亿市场+5用例+协议战国+战略+7层安全+路线图+客户案例）
  - 后半（章节 2/3）保留但语气调整为"抛砖引玉"
  - 总页数: 23

结构:
  PART 0: 封面+目录 (2 页)
  PART 1: AgentCore Payments 深度 (12 页) — 章节分隔 + 11 内容
    1. 章节分隔
    2. 场景钩子
    3. AgentCore 定义
    4. x402 协议
    5. 经济学
    6. 万亿市场 (NEW)
    7. 5 大用例 (NEW)
    8. 协议战国格局 (NEW)
    9. 两大阵营 + AWS 战略 (NEW)
    10. 7 层 Payment Guardrail (NEW)
    11. 路线图三阶段 (NEW)
    12. 客户案例 (NEW)
  PART 2: 短板与转折 (3 页) — 章节分隔 + 短板 + 转折
  PART 3: OpenAgentPay 扩展 (5 页) — 章节分隔 + 5 层架构 + 双协议 + Live + HashKey demo
  PART 4: CTA (1 页)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色
BG       = RGBColor(0x0A, 0x1A, 0x2E)
BG_CARD  = RGBColor(0x14, 0x2A, 0x44)
BG_DARK  = RGBColor(0x05, 0x10, 0x1F)
FG       = RGBColor(0xF0, 0xF4, 0xF8)
DIM      = RGBColor(0x8B, 0xA0, 0xBA)
MUTED    = RGBColor(0x4A, 0x5C, 0x75)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)
GOLD     = RGBColor(0xFF, 0xB8, 0x00)
AWS_ORG  = RGBColor(0xFF, 0x99, 0x00)
OAP_GRN  = RGBColor(0x00, 0xFF, 0x88)
RED      = RGBColor(0xFF, 0x44, 0x44)
PURPLE   = RGBColor(0xBF, 0x5A, 0xF2)

IMG_DIR = Path(__file__).parent / "images"
TOTAL_PAGES = 23

# ============================================================================
#  HELPER FUNCTIONS
# ============================================================================
def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def t(slide, left, top, w, h, text, *, size=18, color=FG, bold=False, align="left", font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb

def mt(slide, left, top, w, h, lines, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb

def bar(slide, left, top, w, h, color=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def card(slide, left, top, w, h, color=BG_CARD, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.adjustments[0] = 0.08
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def hdr(slide, title, subtitle=None, accent=CYAN):
    t(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
      title, size=26, color=FG, bold=True)
    if subtitle:
        t(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
          subtitle, size=13, color=DIM)
    bar(slide, Inches(0.6), Inches(1.4), Inches(2), Emu(40000), accent)

def ftr(slide, page_num=None):
    t(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
      "AgentCore Payments + OpenAgentPay  ·  2026-05-18  ·  Neo Sun",
      size=9, color=MUTED)
    if page_num:
        t(slide, Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.3),
          f"{page_num:02d} / {TOTAL_PAGES:02d}", size=9, color=MUTED, align="right")

def n(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def metric(slide, left, top, w, h, number, label, color=CYAN, num_size=44, lbl_size=11):
    card(slide, left, top, w, h, BG_CARD, color)
    t(slide, left, top + Inches(0.3), w, Inches(1.2),
      number, size=num_size, color=color, bold=True, align="center")
    t(slide, left, top + h - Inches(0.5), w, Inches(0.3),
      label, size=lbl_size, color=DIM, align="center")

def chap(prs, num, title, subtitle, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG_DARK)
    t(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
      f"第 {num} 章", size=24, color=accent, bold=True)
    t(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.5),
      title, size=52, color=FG, bold=True)
    t(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.6),
      subtitle, size=18, color=DIM)
    bar(s, Inches(0.6), Inches(5.4), Inches(3), Emu(60000), accent)
    return s

# ============================================================================
#  BUILD ALL 23 SLIDES
# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ===== Slide 1: 封面（标题改：AgentCore 主，OpenAgentPay 副）=====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG_DARK)
t(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
  "AWS Web3 Team · 周会分享 · 2026-05-18", size=12, color=DIM)
t(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.0),
  "AgentCore Payments", size=58, color=AWS_ORG, bold=True)
t(s, Inches(0.6), Inches(2.95), Inches(12), Inches(0.5),
  "AI Agent 经济基础设施", size=22, color=FG)
t(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
  "兼谈 OpenAgentPay 扩展", size=18, color=OAP_GRN)
t(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
  "对接亚太合规交易所、CEX 钱包、Web3 自托管钱包", size=13, color=DIM)
t(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
  "Neo Sun  ·  AWS Solutions Architect  ·  jiasunm@amazon.com", size=13, color=DIM)
bar(s, Inches(0.6), Inches(6.0), Inches(2), Emu(60000), AWS_ORG)
bar(s, Inches(2.7), Inches(6.0), Inches(2), Emu(60000), OAP_GRN)
n(s, """【封面 · 开场白】

兄弟们好。今天这个分享我准备了挺久的——主题是 AgentCore Payments，附带讲一下我做的 OpenAgentPay 扩展。

简单说一下背景。上周 SSO 给我们团队讲了 AgentCore Payments 这个 AWS 上周（2026-05-07）刚发布的新产品。我对它做了深入研究，写了一篇公众号文章解读它。今天的分享会比公众号文章更聚焦——主要从我们 Web3 团队的角度看：这个产品是什么、它有多重要、它当前覆盖不到的地方在哪里、我们能怎么补。

整个分享分**前后两部分**：

**前面大部分时间（约 10-12 分钟）讲 AgentCore Payments 本身**——这是 AWS 在 Agent 经济卡位的关键基础设施，我会从场景、协议、经济学、用例、协议格局、战略、安全、路线图、客户多个维度讲清楚它。

**后面少部分时间（约 5-8 分钟）讲我做的 OpenAgentPay 扩展**——是个 work-in-progress 的小项目，主要解决 AgentCore Payments 当前对亚太客户覆盖不到的问题。我做这个**主要是想抛砖引玉，邀请兄弟们一起完善**。

兄弟们任何时候**有想法、有挑战、有反对意见都欢迎打断**——今天最理想的状态就是大家七嘴八舌讨论，一起把这个事看清楚。""")
ftr(s, 1)

# ===== Slide 2: 目录 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "今天分享的三件事", "前半 AgentCore Payments 重点讲透  ·  后半 OpenAgentPay 抛砖引玉", AWS_ORG)

# Part 1
card(s, Inches(0.5), Inches(2.0), Inches(8.0), Inches(4.6), BG_CARD, AWS_ORG)
t(s, Inches(0.75), Inches(2.2), Inches(7), Inches(0.5),
  "Part 1  ·  AgentCore Payments", size=22, color=AWS_ORG, bold=True)
t(s, Inches(0.75), Inches(2.7), Inches(7), Inches(0.4),
  "约 10-12 分钟  ·  重点讲透", size=12, color=DIM, bold=True)
mt(s, Inches(0.75), Inches(3.2), Inches(7.5), Inches(3.4), [
    ("· 场景钩子：Agent 为什么需要钱包", 13, FG, False),
    ("· 一句话定义 + 平台地位", 13, FG, False),
    ("· x402 协议怎么工作（5 步）", 13, FG, False),
    ("· 经济学：3000x 优势", 13, FG, False),
    ("· 万亿市场预测（McKinsey）", 13, FG, False),
    ("· 5 大用例 + Heurist 客户", 13, FG, False),
    ("· 2026 协议战国格局（10+ 协议）", 13, FG, False),
    ("· 两大阵营 · AWS 为什么选买方侧", 13, FG, False),
    ("· 7 层 Payment Guardrail（企业级安全）", 13, FG, False),
    ("· 三阶段路线图 + 客户案例", 13, FG, False),
])

# Part 2 + 3
card(s, Inches(8.7), Inches(2.0), Inches(4.2), Inches(2.2), BG_CARD, RED)
t(s, Inches(8.95), Inches(2.2), Inches(4), Inches(0.5),
  "Part 2  ·  短板", size=18, color=RED, bold=True)
t(s, Inches(8.95), Inches(2.7), Inches(4), Inches(0.4),
  "约 2 分钟", size=11, color=DIM)
mt(s, Inches(8.95), Inches(3.1), Inches(4), Inches(1.0), [
    ("· 现状与限制", 12, FG, False),
    ("· 客户在哪？", 12, FG, False),
])

card(s, Inches(8.7), Inches(4.4), Inches(4.2), Inches(2.2), BG_CARD, OAP_GRN)
t(s, Inches(8.95), Inches(4.6), Inches(4), Inches(0.5),
  "Part 3  ·  OpenAgentPay", size=16, color=OAP_GRN, bold=True)
t(s, Inches(8.95), Inches(5.1), Inches(4), Inches(0.4),
  "约 5-8 分钟  ·  抛砖引玉", size=11, color=DIM)
mt(s, Inches(8.95), Inches(5.5), Inches(4), Inches(1.0), [
    ("· 5 层框架 + 双协议轨道", 12, FG, False),
    ("· Live demo + HashKey 实战", 12, FG, False),
    ("· 邀请兄弟们一起完善", 12, FG, False),
])

n(s, """【目录】

今天的分享 3 部分。

**Part 1：AgentCore Payments**——这是今天的重点，约 10-12 分钟。我从 10 个维度把它讲透：场景钩子、平台地位、协议原理、经济学优势、万亿市场预测、5 大官方用例、2026 协议战国格局、两大阵营战略、7 层企业安全设计、三阶段路线图、客户案例。

**Part 2：短板与转折**——快速过，2-3 分钟。AgentCore Payments 是好产品但有短板——只支持 Coinbase + Stripe + 4 个 region（无亚洲）。这是我们 Web3 团队推不动的地方。

**Part 3：OpenAgentPay 扩展**——5-8 分钟。我做的扩展层，已经把 HashKey Chain 接通了。但是 work-in-progress，**主要目的是抛砖引玉，邀请兄弟们一起完善**。

兄弟们听的时候记一个心态——**前半重点听 AgentCore Payments 的产品深度**（这是真的产品级深度研究），**后半听个引子和邀请**（项目本身是小项目，重要的是 community 思路）。

随时打断，最理想的状态是大家七嘴八舌讨论。""")
ftr(s, 2)

# ===== Slide 3: 章节分隔 第 1 章 =====
s = chap(prs, 1, "AgentCore Payments", "AWS 给 AI Agent 发了一张钱包  ·  2026-05-07 Preview 发布", AWS_ORG)
n(s, """【第一章 · AgentCore Payments】

接下来 11 页是今天的重点——AgentCore Payments 深度解读。这 11 页我会从最基础的"为什么要做这个产品"讲到最深度的"AWS 战略和路线图"。

兄弟们听的时候特别注意——这不是简单的产品介绍，是**对 AI Agent 经济卡位战的深度分析**。AgentCore Payments 不只是个新功能，是 AWS 在 Agent 经济基础设施这块的关键卡位。""")
ftr(s, 3)

# ===== Slide 4: 场景钩子 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
t(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
  "想象一下这个场景", size=12, color=DIM)
t(s, Inches(0.6), Inches(1.6), Inches(12), Inches(2.0),
  "你的 AI Agent", size=42, color=FG, bold=True)
t(s, Inches(0.6), Inches(2.5), Inches(12), Inches(2.0),
  "能自己付款吗？", size=72, color=GOLD, bold=True)
t(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
  "在今天之前，答案是：不能。", size=24, color=DIM)
t(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
  "Agent 必须停下来等人类介入：开账号、绑信用卡、管理 API Key、处理订阅", size=14, color=MUTED)
card(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6), BG_CARD, RED)
t(s, Inches(0.8), Inches(6.05), Inches(12), Inches(0.5),
  "智能是自主的。付款，不是。  ←  AI Agent 领域被所有人忽视但致命的矛盾",
  size=14, color=RED, bold=True)
n(s, """【场景钩子】

兄弟们想象一个场景。你的 AI Agent 在帮你做金融研究，分析到一半发现需要访问彭博终端的付费数据才能给出更准确的结论。这个时候 Agent 能自己付款吗？

**在今天之前，答案是不能**。Agent 必须停下来等人类介入：开账号、绑信用卡、管理 API Key、处理订阅...走完一整套人类世界为人类设计的支付流程。

这就是 AI Agent 领域一个被所有人忽视但致命的矛盾——**智能是自主的，付款不是**。

为什么这是个大问题？Agent 经济正在快速到来。Anthropic 的 Computer Use、OpenAI 的 Operator、AWS 自己的 Strands Agent，都在让 Agent 能自主使用工具。但**只要付款这一环还要人参与，Agent 的自主性就被打了折扣**——只能做"调用免费 API"，不能做"研究决策时按需付费查数据"这种真正有商业价值的任务。

AWS 的 AgentCore Payments 就是要解决这个问题——给 AI Agent 一张能自己用的钱包。

**这是 Agent 经济的关键基础设施层**——AWS 卡位的产品。""")
ftr(s, 4)

# ===== Slide 5: AgentCore 一句话定义 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "AgentCore Payments 是什么",
    "Amazon Bedrock AgentCore 平台的原生支付模块  ·  2026-05-07 Preview 发布", AWS_ORG)
card(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(1.0), BG_CARD, AWS_ORG)
t(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.4),
  "首个面向自主 AI Agent 的托管支付基础设施", size=18, color=GOLD, bold=True)
t(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.4),
  "让 Agent 在推理循环内自主发现 → 认证 → 结算 → 获取付费资源，无需人类介入",
  size=12, color=DIM)
img_path = str(IMG_DIR / "01-platform-landscape.png")
s.shapes.add_picture(img_path, Inches(2.0), Inches(2.85), Inches(9.3), Inches(3.7))
t(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
  "三个关键词：首个 (Azure/GCP 都没同类)  ·  托管 (AWS 一手包办)  ·  推理循环内 (不打断 Agent 思考)",
  size=12, color=CYAN, align="center")
n(s, """【AgentCore Payments 是什么】2026-05-07，AWS 在金融服务峰会发布 Amazon Bedrock AgentCore Payments Preview。

一句话定义：**首个面向自主 AI Agent 的托管支付基础设施**——让 Agent 在推理循环内自主发现、认证、结算、获取付费资源，无需人类介入。

注意三个关键词：

**"首个"**——微软 Azure 和 Google Cloud 在 Agent 平台层都**没有同类产品**。AWS 在这个具体赛道**领先一个产品代际**。

**"托管"**——AWS 一手包办协议协商、钱包管理、限额控制、合规检查、可观测性。开发者**几行代码就能接入**。Heurist AI 创始人原话："few lines of code"。

**"推理循环内"**——付款动作不打断 Agent 思考，就像人喘气一样自然。Agent 调用工具的时候碰到 402 Payment Required，AgentCore 自动处理付款，Agent 完全无感。

看这张图，AgentCore Payments 不是独立产品，而是 AgentCore 平台的**原生模块**，跟 Identity、Gateway、Observability、Memory、Runtime 平起平坐。这就是 AWS 的设计哲学——**'native, not bolted-on'，原生集成不是贴片**。意味着企业安全团队审批一次就行。

Payments 模块包含 4 个子能力：Payment Manager（编排大脑）+ Payment Guardrail（风控边界）+ Wallet Providers（钱包对接）+ Protocols（协议层）。""")
ftr(s, 5)

# ===== Slide 6: x402 协议 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "x402 协议  ·  把 HTTP 402 状态码激活了",
    "Coinbase 2024 年底开放协议  ·  5 步搞定一次 Agent 付款", AWS_ORG)
img_path = str(IMG_DIR / "02-x402-sequence.png")
s.shapes.add_picture(img_path, Inches(2.0), Inches(1.7), Inches(9.3), Inches(3.7))
metric(s, Inches(0.6),  Inches(5.7), Inches(2.4), Inches(1.2), "75M",   "近 30 天 tx", CYAN, num_size=32)
metric(s, Inches(3.1),  Inches(5.7), Inches(2.4), Inches(1.2), "169M+", "累计 tx", GOLD, num_size=32)
metric(s, Inches(5.6),  Inches(5.7), Inches(2.4), Inches(1.2), "590K+", "活跃买家", OAP_GRN, num_size=32)
metric(s, Inches(8.1),  Inches(5.7), Inches(2.4), Inches(1.2), "100K+", "活跃卖家", PURPLE, num_size=32)
metric(s, Inches(10.6), Inches(5.7), Inches(2.1), Inches(1.2), "~2s",   "L2 finalize", AWS_ORG, num_size=32)
n(s, """【x402 协议】要理解 AgentCore Payments，必须先懂 x402。x402 是 Coinbase 2024 年底发布的开放协议——把 HTTP 协议里那个一直没人用的状态码 **402 Payment Required** 激活了。

**5 步搞定一次 Agent 付款**：

1. Agent 发起请求 GET /resource，普通 HTTP 调用。
2. 服务端返回 **402 Payment Required**，告诉你要付多少、付给谁、用啥币。
3. Agent 内部，AgentCore Payments 查预算 → 拿钱包私钥 → 做 EIP-3009 离线签名。
4. Agent 带上 **X-Payment** 头重试请求。
5. 服务端验证签名 → Facilitator 上链结算 → 返回 **200 OK** + 你要的内容。

整个流程，**Agent 不需要持有 gas、不需要连区块链节点、不管理 nonce**。复杂度全部被 x402 Facilitator 吸收掉了。Agent 的 UX 跟调 REST API 一模一样。

底部这一组数字告诉兄弟们 **x402 不是玩具**：30 天 7541 万笔交易、累计 1.69 亿笔、活跃买家 59 万、卖家 10 万、L2 结算约 2 秒完成。

被 Stripe、AWS、Cloudflare、Vercel、Quicknode、Messari、Alchemy 都信赖。AWS 选 x402 不是保守的选择——是押**最成熟、已生产验证、增长最快**的赛道。

【插话】其实 2026 上半年是 Agent 支付协议的"战国时代"，10+ 个协议并存。AWS 选 x402 是因为它在生产数据 + 生态成熟度 + 技术成熟度三个维度都最强。这个我们后面有专门一页讲。""")
ftr(s, 6)

# ===== Slide 7: 经济学 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "为什么不能用传统信用卡  ·  经济学不允许",
    "x402 vs Stripe 信用卡  ·  扎心的 3000 倍差距", AWS_ORG)
img_path = str(IMG_DIR / "08-economics.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(7.0), Inches(4.3))

card(s, Inches(7.9), Inches(1.7), Inches(4.8), Inches(2.0), BG_CARD, RED)
t(s, Inches(8.1), Inches(1.85), Inches(4.5), Inches(0.4), "传统 Stripe 信用卡", size=14, color=RED, bold=True)
t(s, Inches(8.1), Inches(2.3), Inches(4.5), Inches(0.8), "$0.30", size=44, color=RED, bold=True)
t(s, Inches(8.1), Inches(3.1), Inches(4.5), Inches(0.5), "单笔最低手续费", size=12, color=DIM)

card(s, Inches(7.9), Inches(3.85), Inches(4.8), Inches(2.0), BG_CARD, OAP_GRN)
t(s, Inches(8.1), Inches(4.0), Inches(4.5), Inches(0.4), "x402 + USDC on Base L2", size=14, color=OAP_GRN, bold=True)
t(s, Inches(8.1), Inches(4.45), Inches(4.5), Inches(0.8), "$0.0001", size=44, color=OAP_GRN, bold=True)
t(s, Inches(8.1), Inches(5.25), Inches(4.5), Inches(0.5), "单笔成本", size=12, color=DIM)

card(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7), BG_CARD, GOLD)
t(s, Inches(0.8), Inches(6.18), Inches(12), Inches(0.55),
  "差了 3000 倍。这就是为什么所有按次计费的 Pay-per-use 市场都上不了规模 — 直到 x402 出现。",
  size=14, color=GOLD, bold=True, align="center")

n(s, """【经济学】兄弟们可能问：让 Agent 自动付款，直接绑张信用卡不就完了？答案是：经济学不允许。

看这张图：**传统 Stripe 信用卡单笔最低手续费 0.30 美元，x402 单笔成本约 0.0001 美元——差了整整 3000 倍**。

具体例子：Agent 调用一次 API 可能只值 0.001 美元（0.1 分钱）。但走传统信用卡要扣 0.30 美元（30 分钱）手续费——**付款比买的东西贵了 300 倍**，这个市场根本上不了规模。

所有按次计费的 Pay-per-use 市场——API 调用、AI 推理调用、按字数计费的内容、按次的数据查询——都被传统支付的单位经济学**彻底锁死**几十年。Agent 经济需要一个新支付层。

而 x402 + USDC on Base L2：
- **L2 结算 ~2 秒**（Flashblock 预确认 200ms），传统 ACH/Card 要 T+2-3 天
- **单笔成本 < 1 分钱**，让小于 1 美元的微支付变得理所当然
- **零账户设置**——有钱包就能付，不用注册商户
- **无 API Key 风险**——签名是一次性授权，泄露也无害
- **天然全球同价**——不管在哪个国家，费率都一样

让小于 1 美元的 Agent 微支付**从不可能变成理所当然**。这就是 AWS、Stripe、Cloudflare 同时押注这个协议的真正原因——**不是因为它酷，而是因为它解锁了一个被锁住几十年的万亿级市场**。""")
ftr(s, 7)

# ===== Slide 8: 万亿市场预测（NEW）=====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "这个市场有多大  ·  万亿美元级别",
    "McKinsey 2030 年预测  ·  Agent 经济正在改变商业根基", GOLD)

# 大数字 hero
t(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
  "McKinsey 预测：到 2030 年", size=18, color=DIM)

t(s, Inches(0.6), Inches(2.7), Inches(12), Inches(2.0),
  "$3T - $5T", size=120, color=GOLD, bold=True, align="center")

t(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
  "Agentic Commerce 中介化的全球商业规模", size=18, color=FG, align="center")

# 关键引用
card(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.4), BG_CARD, CYAN)
t(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(0.4),
  "其中美国 B2C 零售单独就有 $1T 被 Agent 编排", size=15, color=CYAN, bold=True)
t(s, Inches(0.85), Inches(6.0), Inches(11.7), Inches(0.4),
  '"There will soon be more AI agents transacting than humans."', size=12, color=DIM, font="Calibri")
t(s, Inches(0.85), Inches(6.4), Inches(11.7), Inches(0.3),
  "— Brian Foster, Head of Infrastructure Growth, Coinbase", size=10, color=MUTED)

n(s, """【万亿市场】

兄弟们有人会觉得"Agent 自己付款"还是个 niche 场景。但 McKinsey 的预测是**惊人的**：

**到 2030 年，Agentic Commerce 将中介化 $3T 到 $5T 全球商业规模**。其中美国 B2C 零售市场单独就有 **$1T 美元**被 Agent 编排。

什么意思？2030 年时：
- 你在 Amazon 买东西，可能不是你直接搜，是你的 Agent 找最优价格然后下单
- 你订机票，Agent 比价 + 锁价 + 自动付款
- 公司的差旅报销，Agent 自动处理 + 自动支付
- B2B 采购，公司 Agent 跟供应商 Agent 谈判 + 下单 + 付款

**Coinbase Brian Foster 的判断更激进**——"There will soon be more AI agents transacting than humans"——交易的 AI Agent 数量将超过人类。

这是个量级转变——不是"AI 帮人付款"，是"**AI 自己作为经济主体在交易**"。如果你信这个判断，那 AgentCore Payments 就是这个万亿级市场的入口基础设施。

**对我们 Web3 SA 团队的意义**：这是 AWS Agent 经济卡位战的产品级武器。我们的客户（亚洲交易所、Web3 平台）想分这个万亿市场的一杯羹，必须有 Agent Payments 能力。问题是 AgentCore Payments 当前覆盖不到他们——这就是后半 OpenAgentPay 的起点。""")
ftr(s, 8)

# ===== Slide 9: 5 大用例 + Heurist 客户案例 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "5 大官方用例  ·  Heurist 已生产上线",
    "金融分析  ·  内容订阅  ·  数据 API  ·  开发工具  ·  AI 模型调用", AWS_ORG)

img_path = str(IMG_DIR / "03-use-cases.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(7.5), Inches(4.5))

# 右半：Heurist 客户案例引用
card(s, Inches(8.4), Inches(1.7), Inches(4.4), Inches(4.5), BG_CARD, OAP_GRN)
t(s, Inches(8.6), Inches(1.85), Inches(4.0), Inches(0.4),
  "Heurist AI", size=18, color=OAP_GRN, bold=True)
t(s, Inches(8.6), Inches(2.3), Inches(4.0), Inches(0.4),
  "已生产上线 ✓", size=12, color=OAP_GRN)

t(s, Inches(8.6), Inches(2.85), Inches(4.0), Inches(0.4),
  "JW Wang，创始人", size=11, color=DIM, bold=True)

mt(s, Inches(8.6), Inches(3.25), Inches(4.0), Inches(2.4), [
    ("\"Heurist is using AgentCore", 11, FG, False),
    ("payments for our research", 11, FG, False),
    ("agent. End customers can set", 11, FG, False),
    ("a budget for the research and", 11, FG, False),
    ("the agent uses AgentCore", 11, FG, False),
    ("payments to get accurate", 11, FG, False),
    ("real-time data.", 11, FG, False),
    ("", 8, DIM, False),
    ("We were able to integrate", 11, FG, False),
    ("payments quickly with low", 11, FG, False),
    ("effort and few lines of code.\"", 11, GOLD, True),
])

n(s, """【5 大官方用例】

AWS 文档里权威列举了 5 大核心用例。看这张图——

**5 大用例**：
1. **金融分析**——Agent 按需购买实时行情、社交舆情、新闻数据（Heurist 在做）
2. **内容订阅**——Agent 按需付费访问付费墙后的内容（Warner Bros 在评估）
3. **数据 API**——Agent 按调用计费访问数据服务（彭博、Refinitiv 这类）
4. **开发工具**——Agent 按 token 计费用 SaaS（GitHub Copilot 类）
5. **AI 模型调用**——Agent 调用其他 Agent 服务（agent-to-agent 经济）

外加 **4 个 FSI 行业扩展场景**（图中下半部分）。

**右边这个引用是关键证据**——Heurist AI 创始人 JW Wang 在 AWS 官方 Blog 里直接背书：

> "Heurist 在用 AgentCore Payments 做研究 Agent。终端客户可以给研究设预算，Agent 自动用 AgentCore Payments 购买实时数据。我们用**很少几行代码**就接进去了。"

注意那句 **"few lines of code"**——这是 AgentCore Payments 整个产品最核心的卖点：**把几个月的工程投入降到几行代码**。

为什么这个引用重要？因为：
- **Heurist 不是 PR 案例，是 production**——已生产上线
- **客户现身说法**比任何官方文案都有力
- 证明产品**真的落地了**——不是 Vapor ware

Heurist 是当前**最重要的 AgentCore Payments 客户案例**——任何兄弟讲 AgentCore Payments 都该提这个。""")
ftr(s, 9)

# ===== Slide 10: 协议战国格局（NEW）=====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG)
hdr(s, "2026 协议战国格局  ·  10+ 协议并存零互操作",
    "Q2 2026 行业现实  ·  AWS 在多个协议中押了最成熟那一个", PURPLE)

img_path = str(IMG_DIR / "05-protocol-war.png")
s.shapes.add_picture(img_path, Inches(0.6), Inches(1.7), Inches(8.5), Inches(5.0))

# 右半：协议表格
card(s, Inches(9.3), Inches(1.7), Inches(3.6), Inches(5.0), BG_CARD, PURPLE)
t(s, Inches(9.5), Inches(1.85), Inches(3.3), Inches(0.4),
  "主要协议 + 推手", size=13, color=PURPLE, bold=True)

mt(s, Inches(9.5), Inches(2.3), Inches(3.3), Inches(4.2), [
    ("买方侧 (Agent 发起)", 11, CYAN, True),
    ("· x402 ★ Coinbase → AWS", 10, FG, False),
    ("· MPP   Stripe + Tempo", 10, DIM, False),
    ("· ACP   OpenAI", 10, DIM, False),
    ("", 8, DIM, False),
    ("收款侧 (商户接受)", 11, GOLD, True),
    ("· AP2   Google → FIDO", 10, DIM, False),
    ("· UCP   Google + Visa + MC", 10, DIM, False),
    ("· TAP   Visa", 10, DIM, False),
    ("· Agent Pay   Mastercard", 10, DIM, False),
    ("", 8, DIM, False),
    ("AWS 选 x402 因为：", 11, AWS_ORG, True),
    ("· 生产验证 (75M tx/月)", 10, FG, False),
    ("· 生态成熟 (Stripe 也用)", 10, FG, False),
    ("· 技术成熟 (1 年 169M tx)", 10, FG, False),
])

n(s, """【协议战国格局】

兄弟们可能好奇：只有 AWS 在做这件事吗？

**恰恰相反**——2026 年上半年是 AI Agent 支付协议的"战国时代"。几乎每家主流支付网络和科技巨头都甩出了自己的协议。

看这张图——**10+ 个协议并存，零互操作性**（这是 Q2 2026 digitalapplied.com 明确报告的行业现实）。

我帮你梳理最关键的几家：

**Agent 买方侧**（让 Agent 能发起付款）：
- **x402** ★ — Coinbase 推，AWS 首选
- **MPP** — Stripe + Tempo (IETF Internet-Draft，向后兼容 x402)
- **ACP** — OpenAI (ChatGPT 原生集成)

**Merchant 收款侧**（让商户能接受 Agent 支付）：
- **AP2** — Google 推，捐 FIDO Alliance
- **UCP** — Google (Shopify + Walmart + Visa + Mastercard 背书)
- **TAP** — Visa (Trusted Agent Protocol)
- **Agent Pay** — Mastercard (集成 Microsoft Copilot)

**关键判断**：Azure 和 GCP 在 Agent **平台层都没有**同类 Payments 原生产品——AWS 在这个具体赛道**领先一个产品代际**。

**AWS 为什么选 x402**？三个理由（右边卡片）：
- 生产验证 (75M tx/月，被 Stripe/Cloudflare 信赖)
- 生态成熟 (10K+ x402 endpoints in Bazaar)
- 技术成熟 (1 年累计 169M tx，已稳定 1 年)

不是因为它酷，是因为它**最经过考验**。

**对我们 SA 的意义**：客户问"为什么 AWS 选 x402 不选 MPP"——答案是技术成熟度。MPP 是 Stripe + Tempo 的 IETF 草案，2025 才出，未经生产验证。x402 已经 169M tx 跑过。""")
ftr(s, 10)

print(f"Done with first half (slides 1-10)")
print(f"Total slides so far: {len(prs.slides)}")

# 保存中间状态
out = Path(__file__).parent / "openagentpay-talk-half.pptx"
prs.save(out)
print(f"✅ Saved intermediate: {out}")
